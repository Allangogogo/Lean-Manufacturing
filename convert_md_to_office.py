#!/usr/bin/env python3
"""
convert_md_to_office.py
=======================
Converts Markdown files in the Lean knowledge base to Office formats
(Word .docx, Excel .xlsx, PowerPoint .pptx) based on directory structure.

Usage:
    python convert_md_to_office.py [--dry-run] [--keep-md]

Directory routing:
    - 01-精益工具知识库/ (all subdirectories)  -> .docx
    - 02-精益培训/01-培训策略/                   -> .docx
    - 02-精益培训/02-培训材料/                   -> .docx (4 files -> .pptx)
    - 02-精益培训/03-培训计划/                   -> .xlsx
    - 02-精益培训/04-培训模板/                   -> .xlsx
    - 02-精益培训/05-记录追踪/                   -> .xlsx
    - 02-精益培训/06-效果反馈/                   -> .docx
    - 03-成熟度评估/01-评估框架/                 -> .docx
    - 03-成熟度评估/02-工厂整体评估/             -> .docx
    - 03-成熟度评估/03-局部评估/                 -> .xlsx
    - 03-成熟度评估/04-评估报告/                 -> .docx
    - 04-实施战略/01-实施战略/                   -> .docx
    - 04-实施战略/02-详细计划/                   -> .docx
    - 04-实施战略/03-实施工具/                   -> .xlsx
    - 05-项目管理/01-项目治理/                   -> .docx
    - 05-项目管理/02-进度管理/                   -> .xlsx
    - 05-项目管理/03-风险管理/                   -> .docx
    - 05-项目管理/04-绩效管理/                   -> .docx
    - appendix/                                  -> .docx
"""

import os
import re
import sys
import argparse
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any

# ---------------------------------------------------------------------------
# Third-party imports
# ---------------------------------------------------------------------------
try:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
except ImportError:
    sys.exit("ERROR: python-docx is required. Install with: pip install python-docx")

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
except ImportError:
    sys.exit("ERROR: openpyxl is required. Install with: pip install openpyxl")

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("ERROR: python-pptx is required. Install with: pip install python-pptx")


# ===========================================================================
# Configuration
# ===========================================================================

BASE_DIR = Path(__file__).resolve().parent

# Files to always skip (keep as .md)
SKIP_FILES = {"README.md", "CLAUDE.md"}

# ---- Format routing ----
# Directories that produce .xlsx
XLSX_DIRS = {
    "02-精益培训/03-培训计划",
    "02-精益培训/04-培训模板",
    "02-精益培训/05-记录追踪",
    "03-成熟度评估/03-局部评估",
    "04-实施战略/03-实施工具",
    "05-项目管理/02-进度管理",
}

# Specific files (relative to BASE_DIR) that produce .pptx
PPTX_FILES = {
    "02-精益培训/02-培训材料/01-精益意识培训大纲.md",
    "02-精益培训/02-培训材料/02-工具专项培训大纲.md",
    "02-精益培训/02-培训材料/03-领导力培训大纲.md",
    "02-精益培训/02-培训材料/04-高级精益培训大纲.md",
}

# Everything else (non-skip, non-xlsx, non-pptx) goes to .docx.

# ---- Font configuration ----
FONT_BODY_CN = "SimSun"       # Chinese body text
FONT_HEADING_CN = "SimHei"    # Chinese headings
FONT_CODE = "Courier New"     # Code blocks
FONT_BODY_SIZE = Pt(11)
FONT_HEADING1_SIZE = Pt(22)
FONT_HEADING2_SIZE = Pt(16)
FONT_HEADING3_SIZE = Pt(13)
FONT_CODE_SIZE = Pt(9)

# ---- Color palette ----
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
TABLE_HEADER_BG = "1F3A5F"
TABLE_ALT_ROW_BG = "EDF2F9"


# ===========================================================================
# Markdown Parser — Intermediate Representation
# ===========================================================================
# We parse markdown into a list of "blocks". Each block is a dict with a
# 'type' key and relevant data. This lets each converter consume the same
# parsed structure.

def parse_markdown(text: str) -> List[Dict[str, Any]]:
    """
    Parse markdown text into a list of block dicts.

    Block types:
        heading      — {level: 1|2|3|4, text: str}
        paragraph    — {text: str}
        table        — {headers: [str], rows: [[str]]}
        bullet_list  — {items: [str]}
        numbered_list — {items: [str]}
        code_block   — {lang: str, code: str}
        horizontal_rule — {}
        blockquote   — {text: str}
    """
    lines = text.split("\n")
    blocks: List[Dict[str, Any]] = []
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # --- Empty line → skip ---
        if not stripped:
            i += 1
            continue

        # --- Horizontal rule ---
        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", stripped):
            blocks.append({"type": "horizontal_rule"})
            i += 1
            continue

        # --- Heading ---
        heading_match = re.match(r"^(#{1,6})\s+(.*)", stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 4)
            blocks.append({"type": "heading", "level": level,
                           "text": heading_match.group(2).strip()})
            i += 1
            continue

        # --- Code block (fenced) ---
        if stripped.startswith("```"):
            lang = stripped[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            blocks.append({"type": "code_block", "lang": lang,
                           "code": "\n".join(code_lines)})
            i += 1  # skip closing ```
            continue

        # --- Table ---
        if "|" in stripped and stripped.startswith("|"):
            table_lines = []
            while i < len(lines) and "|" in lines[i].strip() and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            blocks.append(_parse_table(table_lines))
            continue

        # --- Blockquote ---
        if stripped.startswith(">"):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(
                    re.sub(r"^>\s?", "", lines[i].strip()))
                i += 1
            blocks.append({"type": "blockquote",
                           "text": " ".join(quote_lines)})
            continue

        # --- Unordered list ---
        if re.match(r"^[-*+]\s+", stripped):
            items = []
            while i < len(lines) and re.match(
                    r"^\s*[-*+]\s+", lines[i]):
                item_text = re.sub(r"^\s*[-*+]\s+", "", lines[i])
                items.append(item_text.strip())
                i += 1
            blocks.append({"type": "bullet_list", "items": items})
            continue

        # --- Ordered list ---
        if re.match(r"^\d+[.)]\s+", stripped):
            items = []
            while i < len(lines) and re.match(
                    r"^\s*\d+[.)]\s+", lines[i]):
                item_text = re.sub(r"^\s*\d+[.)]\s+", "", lines[i])
                items.append(item_text.strip())
                i += 1
            blocks.append({"type": "numbered_list", "items": items})
            continue

        # --- Paragraph (default) — accumulate contiguous non-empty lines ---
        para_lines = []
        while i < len(lines) and lines[i].strip() and \
                not _is_block_start(lines[i]):
            para_lines.append(lines[i].strip())
            i += 1
        if para_lines:
            blocks.append({"type": "paragraph",
                           "text": " ".join(para_lines)})

    return blocks


def _is_block_start(line: str) -> bool:
    """Return True if the line starts a new block type."""
    s = line.strip()
    if not s:
        return True
    if s.startswith("#"):
        return True
    if s.startswith("```"):
        return True
    if re.match(r"^(-{3,}|\*{3,}|_{3,})$", s):
        return True
    if s.startswith("|"):
        return True
    if s.startswith(">"):
        return True
    if re.match(r"^[-*+]\s+", s):
        return True
    if re.match(r"^\d+[.)]\s+", s):
        return True
    return False


def _parse_table(lines: List[str]) -> Dict[str, Any]:
    """Parse markdown table lines into headers + rows."""
    if len(lines) < 2:
        return {"type": "paragraph", "text": "\n".join(lines)}

    def split_row(line: str) -> List[str]:
        """Split a markdown table row by | and strip cells."""
        cells = line.split("|")
        # Remove leading/trailing empty cells from the split
        if cells and cells[0].strip() == "":
            cells = cells[1:]
        if cells and cells[-1].strip() == "":
            cells = cells[:-1]
        return [c.strip() for c in cells]

    headers = split_row(lines[0])
    # Line 1 is the separator (|---|---|), skip it
    rows = []
    for line in lines[2:]:
        # Skip separator lines
        if re.match(r"^\|[\s\-:|]+\|$", line):
            continue
        rows.append(split_row(line))

    return {"type": "table", "headers": headers, "rows": rows}


# ===========================================================================
# Inline Markdown Parser (bold, italic, code, links)
# ===========================================================================

# Pattern for inline elements: **bold**, *italic*, `code`, [text](url)
_INLINE_PATTERN = re.compile(
    r"(\*\*(.+?)\*\*)"        # bold
    r"|(\*(.+?)\*)"           # italic
    r"|(`([^`]+)`)"           # inline code
    r"|(\[([^\]]+)\]\([^)]+\))"  # link — we extract text only
)


def _strip_inline_md(text: str) -> str:
    """Remove markdown inline formatting, returning plain text."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text


# ===========================================================================
# Word (.docx) Converter
# ===========================================================================

def _set_cell_border(cell, **kwargs):
    """
    Set cell border. Usage:
        _set_cell_border(cell, top={"sz": 6, "color": "000000", "val": "single"})
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = parse_xml(f'<w:tcBorders {nsdecls("w")}></w:tcBorders>')
    for edge, attrs in kwargs.items():
        element = parse_xml(
            f'<w:{edge} {nsdecls("w")} w:val="{attrs["val"]}" '
            f'w:sz="{attrs["sz"]}" w:space="0" '
            f'w:color="{attrs["color"]}"/>')
        tcBorders.append(element)
    tcPr.append(tcBorders)


def _add_formatted_runs(paragraph, text: str, base_font_name: str = FONT_BODY_CN,
                        base_font_size: Pt = FONT_BODY_SIZE,
                        base_bold: bool = False, base_italic: bool = False):
    """Add runs to a paragraph, parsing inline markdown formatting."""
    pos = 0
    for match in _INLINE_PATTERN.finditer(text):
        # Add plain text before this match
        if match.start() > pos:
            plain = text[pos:match.start()]
            run = paragraph.add_run(plain)
            run.font.name = base_font_name
            run.font.size = base_font_size
            run.bold = base_bold
            run.italic = base_italic

        # Bold
        if match.group(2):
            run = paragraph.add_run(match.group(2))
            run.font.name = base_font_name
            run.font.size = base_font_size
            run.bold = True
            run.italic = base_italic
        # Italic
        elif match.group(4):
            run = paragraph.add_run(match.group(4))
            run.font.name = base_font_name
            run.font.size = base_font_size
            run.bold = base_bold
            run.italic = True
        # Inline code
        elif match.group(6):
            run = paragraph.add_run(match.group(6))
            run.font.name = FONT_CODE
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0xC7, 0x25, 0x4E)
        # Link (just text)
        elif match.group(8):
            run = paragraph.add_run(match.group(8))
            run.font.name = base_font_name
            run.font.size = base_font_size
            run.font.color.rgb = MEDIUM_BLUE

        pos = match.end()

    # Remaining text after last match
    if pos < len(text):
        run = paragraph.add_run(text[pos:])
        run.font.name = base_font_name
        run.font.size = base_font_size
        run.bold = base_bold
        run.italic = base_italic


def _set_run_font(run, font_name: str, font_size: Pt,
                  bold: bool = False, italic: bool = False):
    """Configure a run's font properties with East Asian font fallback."""
    run.font.name = font_name
    run.font.size = font_size
    run.bold = bold
    run.italic = italic
    # Set East Asian font for Chinese characters
    r = run._element
    rPr = r.find(qn('w:rPr'))
    if rPr is None:
        rPr = parse_xml(f'<w:rPr {nsdecls("w")}></w:rPr>')
        r.insert(0, rPr)
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = parse_xml(f'<w:rFonts {nsdecls("w")}/>')
        rPr.insert(0, rFonts)
    rFonts.set(qn('w:eastAsia'), font_name)


def convert_to_docx(md_path: Path, output_path: Path) -> None:
    """Convert a markdown file to a Word .docx document."""
    text = md_path.read_text(encoding="utf-8")
    blocks = parse_markdown(text)

    doc = Document()

    # -- Set default document style --
    style = doc.styles["Normal"]
    style.font.name = FONT_BODY_CN
    style.font.size = FONT_BODY_SIZE
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.15
    # Set East Asian font on the style
    rPr = style.element.find(qn('w:rPr'))
    if rPr is None:
        rPr = parse_xml(f'<w:rPr {nsdecls("w")}></w:rPr>')
        style.element.append(rPr)
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = parse_xml(f'<w:rFonts {nsdecls("w")}/>')
        rPr.insert(0, rFonts)
    rFonts.set(qn('w:eastAsia'), FONT_BODY_CN)

    # Configure heading styles
    for lvl, size in [(1, FONT_HEADING1_SIZE), (2, FONT_HEADING2_SIZE),
                      (3, FONT_HEADING3_SIZE), (4, Pt(12))]:
        style_name = f"Heading {lvl}"
        if style_name in doc.styles:
            hs = doc.styles[style_name]
            hs.font.name = FONT_HEADING_CN
            hs.font.size = size
            hs.font.color.rgb = DARK_BLUE
            hs.font.bold = True
            hPr = hs.element.find(qn('w:rPr'))
            if hPr is None:
                hPr = parse_xml(f'<w:rPr {nsdecls("w")}></w:rPr>')
                hs.element.append(hPr)
            hFonts = hPr.find(qn('w:rFonts'))
            if hFonts is None:
                hFonts = parse_xml(f'<w:rFonts {nsdecls("w")}/>')
                hPr.insert(0, hFonts)
            hFonts.set(qn('w:eastAsia'), FONT_HEADING_CN)

    # -- Set page margins --
    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(3.18)
        section.right_margin = Cm(3.18)

    # -- Process blocks --
    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = block["level"]
            para = doc.add_heading(level=level)
            _add_formatted_runs(para, block["text"],
                                base_font_name=FONT_HEADING_CN,
                                base_font_size=[
                                    FONT_HEADING1_SIZE, FONT_HEADING2_SIZE,
                                    FONT_HEADING3_SIZE, Pt(12)][min(level - 1, 3)],
                                base_bold=True)

        elif btype == "paragraph":
            para = doc.add_paragraph()
            para.paragraph_format.space_after = Pt(6)
            _add_formatted_runs(para, block["text"])

        elif btype == "table":
            _add_docx_table(doc, block)

        elif btype == "bullet_list":
            for item in block["items"]:
                para = doc.add_paragraph(style="List Bullet")
                _add_formatted_runs(para, item)

        elif btype == "numbered_list":
            for item in block["items"]:
                para = doc.add_paragraph(style="List Number")
                _add_formatted_runs(para, item)

        elif btype == "code_block":
            _add_docx_code_block(doc, block["code"])

        elif btype == "blockquote":
            para = doc.add_paragraph()
            para.paragraph_format.left_indent = Cm(1.0)
            _add_formatted_runs(para, block["text"],
                                base_italic=True,
                                base_font_size=Pt(10))
            # Add a vertical bar effect via left border
            pPr = para._element.get_or_add_pPr()
            pBdr = parse_xml(
                f'<w:pBdr {nsdecls("w")}>'
                f'<w:left w:val="single" w:sz="12" w:space="4" '
                f'w:color="2E75B6"/>'
                f'</w:pBdr>')
            pPr.append(pBdr)

        elif btype == "horizontal_rule":
            para = doc.add_paragraph()
            para.paragraph_format.space_before = Pt(6)
            para.paragraph_format.space_after = Pt(6)
            # Add a horizontal line via bottom border
            pPr = para._element.get_or_add_pPr()
            pBdr = parse_xml(
                f'<w:pBdr {nsdecls("w")}>'
                f'<w:bottom w:val="single" w:sz="6" w:space="1" '
                f'w:color="CCCCCC"/>'
                f'</w:pBdr>')
            pPr.append(pBdr)

    # -- Set document title from first heading --
    first_heading = next(
        (b["text"] for b in blocks if b["type"] == "heading"), None)
    if first_heading:
        doc.core_properties.title = first_heading

    doc.save(str(output_path))


def _add_docx_table(doc: Document, block: Dict[str, Any]) -> None:
    """Add a formatted table to the Word document."""
    headers = block["headers"]
    rows = block["rows"]
    num_cols = len(headers)
    # Ensure all rows have same number of columns
    normalized_rows = []
    for row in rows:
        padded = row + [""] * (num_cols - len(row))
        normalized_rows.append(padded[:num_cols])

    table = doc.add_table(rows=1 + len(normalized_rows), cols=num_cols)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    # Header row
    for j, header_text in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = ""
        para = cell.paragraphs[0]
        _add_formatted_runs(para, _strip_inline_md(header_text),
                            base_font_name=FONT_HEADING_CN,
                            base_font_size=Pt(10), base_bold=True)
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        # Header background color
        shading = parse_xml(
            f'<w:shd {nsdecls("w")} w:fill="{TABLE_HEADER_BG}" '
            f'w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(shading)
        # White font for header
        for run in para.runs:
            run.font.color.rgb = WHITE

    # Data rows
    for i, row_data in enumerate(normalized_rows):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i + 1].cells[j]
            cell.text = ""
            para = cell.paragraphs[0]
            _add_formatted_runs(para, _strip_inline_md(cell_text),
                                base_font_size=Pt(10))
            # Alternating row shading
            if i % 2 == 1:
                shading = parse_xml(
                    f'<w:shd {nsdecls("w")} w:fill="{TABLE_ALT_ROW_BG}" '
                    f'w:val="clear"/>')
                cell._tc.get_or_add_tcPr().append(shading)

    # Set column widths proportionally
    if num_cols > 0:
        avail_width = Cm(14.64)  # A4 width minus margins
        col_width = avail_width // num_cols
        for row in table.rows:
            for cell in row.cells:
                cell.width = col_width

    # Add spacing after table
    doc.add_paragraph()


def _add_docx_code_block(doc: Document, code: str) -> None:
    """Add a code block with monospace font and grey background."""
    para = doc.add_paragraph()
    para.paragraph_format.space_before = Pt(6)
    para.paragraph_format.space_after = Pt(6)
    para.paragraph_format.left_indent = Cm(0.5)

    # Grey background via shading
    pPr = para._element.get_or_add_pPr()
    shading = parse_xml(
        f'<w:shd {nsdecls("w")} w:fill="F5F5F5" w:val="clear"/>')
    pPr.append(shading)
    # Add a thin border
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'<w:top w:val="single" w:sz="4" w:space="1" w:color="DDDDDD"/>'
        f'<w:left w:val="single" w:sz="4" w:space="4" w:color="DDDDDD"/>'
        f'<w:bottom w:val="single" w:sz="4" w:space="1" w:color="DDDDDD"/>'
        f'<w:right w:val="single" w:sz="4" w:space="4" w:color="DDDDDD"/>'
        f'</w:pBdr>')
    pPr.append(pBdr)

    run = para.add_run(code)
    run.font.name = FONT_CODE
    run.font.size = FONT_CODE_SIZE
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# ===========================================================================
# Excel (.xlsx) Converter
# ===========================================================================

def convert_to_xlsx(md_path: Path, output_path: Path) -> None:
    """Convert a markdown file to an Excel .xlsx workbook."""
    text = md_path.read_text(encoding="utf-8")
    blocks = parse_markdown(text)

    wb = Workbook()

    # ---- Sheet 1: Full markdown content as "说明" sheet ----
    ws_info = wb.active
    ws_info.title = "说明"
    _write_info_sheet(ws_info, blocks, md_path.stem)

    # ---- Sheet 2+: Extract tables into their own sheets ----
    table_blocks = [b for b in blocks if b["type"] == "table"]
    if table_blocks:
        for idx, table_block in enumerate(table_blocks):
            sheet_name = f"表格{idx + 1}"
            if idx == 0 and len(table_blocks) == 1:
                sheet_name = "数据"
            ws = wb.create_sheet(title=sheet_name[:31])  # Excel limit 31 chars
            _write_table_sheet(ws, table_block)

    wb.save(str(output_path))


def _write_info_sheet(ws, blocks: List[Dict], title: str) -> None:
    """Write the full markdown content into an instruction sheet."""
    ws.sheet_properties.tabColor = "1F3A5F"

    # Title row
    ws.merge_cells("A1:F1")
    title_cell = ws["A1"]
    title_cell.value = title
    title_cell.font = Font(name=FONT_HEADING_CN, size=16, bold=True,
                           color="1F3A5F")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 35

    # Header border
    header_border = Border(
        bottom=Side(style="medium", color="1F3A5F"))
    ws["A1"].border = header_border

    current_row = 3

    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = block["level"]
            sizes = {1: 14, 2: 12, 3: 11, 4: 10}
            ws.merge_cells(f"A{current_row}:F{current_row}")
            cell = ws.cell(row=current_row, column=1,
                           value=_strip_inline_md(block["text"]))
            cell.font = Font(name=FONT_HEADING_CN,
                             size=sizes.get(level, 10),
                             bold=True, color="1F3A5F")
            if level == 1:
                ws.row_dimensions[current_row].height = 25
            current_row += 1

        elif btype == "paragraph":
            ws.merge_cells(f"A{current_row}:F{current_row}")
            cell = ws.cell(row=current_row, column=1,
                           value=_strip_inline_md(block["text"]))
            cell.font = Font(name=FONT_BODY_CN, size=10)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            current_row += 1

        elif btype == "bullet_list":
            for item in block["items"]:
                ws.merge_cells(f"A{current_row}:F{current_row}")
                cell = ws.cell(row=current_row, column=1,
                               value=f"  *  {_strip_inline_md(item)}")
                cell.font = Font(name=FONT_BODY_CN, size=10)
                cell.alignment = Alignment(wrap_text=True)
                current_row += 1

        elif btype == "numbered_list":
            for idx, item in enumerate(block["items"], 1):
                ws.merge_cells(f"A{current_row}:F{current_row}")
                cell = ws.cell(row=current_row, column=1,
                               value=f"  {idx}. {_strip_inline_md(item)}")
                cell.font = Font(name=FONT_BODY_CN, size=10)
                cell.alignment = Alignment(wrap_text=True)
                current_row += 1

        elif btype == "code_block":
            for code_line in block["code"].split("\n"):
                ws.merge_cells(f"A{current_row}:F{current_row}")
                cell = ws.cell(row=current_row, column=1, value=code_line)
                cell.font = Font(name=FONT_CODE, size=9,
                                 color="333333")
                cell.fill = PatternFill(start_color="F5F5F5",
                                        end_color="F5F5F5",
                                        fill_type="solid")
                current_row += 1

        elif btype == "horizontal_rule":
            current_row += 1  # blank row

        elif btype == "table":
            # Show a preview in the info sheet
            ws.merge_cells(f"A{current_row}:F{current_row}")
            cell = ws.cell(
                row=current_row, column=1,
                value=f"[表格数据 - 请查看对应的表格Sheet]")
            cell.font = Font(name=FONT_BODY_CN, size=10,
                             italic=True, color="2E75B6")
            current_row += 1

        # Add spacing
        current_row += 1

    # Set column widths
    ws.column_dimensions["A"].width = 20
    ws.column_dimensions["B"].width = 20
    ws.column_dimensions["C"].width = 20
    ws.column_dimensions["D"].width = 20
    ws.column_dimensions["E"].width = 20
    ws.column_dimensions["F"].width = 20


def _write_table_sheet(ws, table_block: Dict[str, Any]) -> None:
    """Write a parsed markdown table into an Excel sheet with formatting."""
    headers = table_block["headers"]
    rows = table_block["rows"]
    num_cols = len(headers)

    ws.sheet_properties.tabColor = "2E75B6"

    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )

    header_font = Font(name=FONT_HEADING_CN, size=10, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color=TABLE_HEADER_BG,
                              end_color=TABLE_HEADER_BG,
                              fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center",
                                 wrap_text=True)

    body_font = Font(name=FONT_BODY_CN, size=10)
    body_alignment = Alignment(vertical="center", wrap_text=True)
    alt_fill = PatternFill(start_color=TABLE_ALT_ROW_BG,
                           end_color=TABLE_ALT_ROW_BG,
                           fill_type="solid")

    # Write headers
    for j, header_text in enumerate(headers, 1):
        cell = ws.cell(row=1, column=j,
                       value=_strip_inline_md(header_text))
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_alignment
        cell.border = thin_border

    # Write data rows
    for i, row_data in enumerate(rows, 2):
        for j, cell_text in enumerate(row_data, 1):
            if j > num_cols:
                break
            cell = ws.cell(row=i, column=j,
                           value=_strip_inline_md(cell_text))
            cell.font = body_font
            cell.alignment = body_alignment
            cell.border = thin_border
            if (i - 2) % 2 == 1:
                cell.fill = alt_fill

    # Auto-size columns (approximate)
    for j in range(1, num_cols + 1):
        max_len = max(
            (len(str(ws.cell(row=r, column=j).value or ""))
             for r in range(1, len(rows) + 2)),
            default=10
        )
        ws.column_dimensions[ws.cell(row=1, column=j).column_letter].width = \
            min(max(max_len + 4, 10), 40)

    # Freeze header row
    ws.freeze_panes = "A2"


# ===========================================================================
# PowerPoint (.pptx) Converter
# ===========================================================================

def _pptx_set_font(run, font_name: str = FONT_HEADING_CN,
                   font_size: int = 18, bold: bool = False,
                   color: PptxRGBColor = PptxRGBColor(0, 0, 0)):
    """Configure font properties on a pptx text run."""
    run.font.name = font_name
    run.font.size = PptxPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def convert_to_pptx(md_path: Path, output_path: Path) -> None:
    """Convert a markdown file to a PowerPoint .pptx presentation."""
    text = md_path.read_text(encoding="utf-8")
    blocks = parse_markdown(text)

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)   # 16:9 widescreen
    prs.slide_height = PptxInches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # -- Color constants --
    DARK_BG = PptxRGBColor(0x1F, 0x3A, 0x5F)
    ACCENT = PptxRGBColor(0x2E, 0x75, 0xB6)
    TEXT_DARK = PptxRGBColor(0x33, 0x33, 0x33)
    TEXT_WHITE = PptxRGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = PptxRGBColor(0xF5, 0xF7, 0xFA)

    # -- Extract document title --
    doc_title = md_path.stem
    for b in blocks:
        if b["type"] == "heading" and b["level"] == 1:
            doc_title = _strip_inline_md(b["text"])
            break

    # ===== Title Slide =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Dark blue background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title text box
    txBox = slide.shapes.add_textbox(
        PptxInches(1.0), PptxInches(2.0),
        PptxInches(11.333), PptxInches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = doc_title
    p.alignment = PP_ALIGN.CENTER
    _pptx_set_font(p.runs[0], FONT_HEADING_CN, 36, True, TEXT_WHITE)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = f"精益知识库 - {md_path.parent.name}"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = PptxPt(20)
    run2 = p2.add_run()
    _pptx_set_font(run2, FONT_BODY_CN, 18, False,
                   PptxRGBColor(0xAA, 0xCC, 0xEE))

    # ===== Content Slides =====
    # Group blocks by ## headings
    current_section_title = None
    current_section_blocks: List[Dict] = []

    def _flush_section():
        """Create a slide for the accumulated section blocks."""
        if current_section_title is None and not current_section_blocks:
            return
        _create_content_slide(prs, current_section_title,
                              current_section_blocks)

    for block in blocks:
        btype = block["type"]

        if btype == "heading" and block["level"] == 2:
            # Flush previous section
            if current_section_title is not None or current_section_blocks:
                _flush_section()
            current_section_title = block["text"]
            current_section_blocks = []
        elif btype == "heading" and block["level"] == 1:
            # Skip the top-level heading (already in title slide)
            continue
        else:
            current_section_blocks.append(block)

    # Flush last section
    _flush_section()

    # -- Add slide numbers --
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue  # Skip title slide
        txBox = slide.shapes.add_textbox(
            PptxInches(12.0), PptxInches(7.0),
            PptxInches(1.0), PptxInches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.alignment = PP_ALIGN.RIGHT
        _pptx_set_font(p.add_run(), FONT_BODY_CN, 10, False,
                       PptxRGBColor(0x99, 0x99, 0x99))

    prs.save(str(output_path))


def _create_content_slide(prs: Presentation,
                          title: Optional[str],
                          blocks: List[Dict]) -> None:
    """Create a single content slide from a list of blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    DARK_BG = PptxRGBColor(0x1F, 0x3A, 0x5F)
    ACCENT = PptxRGBColor(0x2E, 0x75, 0xB6)
    TEXT_DARK = PptxRGBColor(0x33, 0x33, 0x33)
    LIGHT_BG = PptxRGBColor(0xF5, 0xF7, 0xFA)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # -- Title bar (dark blue strip at top) --
    title_bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        PptxInches(0), PptxInches(0),
        slide_width, PptxInches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BG
    title_bar.line.fill.background()

    if title:
        txBox = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.15),
            PptxInches(12.0), PptxInches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = _strip_inline_md(title)
        _pptx_set_font(p.runs[0], FONT_HEADING_CN, 28, True,
                       PptxRGBColor(0xFF, 0xFF, 0xFF))

    # -- Content area --
    content_top = PptxInches(1.3)
    content_height = PptxInches(5.8)
    content_left = PptxInches(0.6)
    content_width = PptxInches(12.1)

    # Determine content type and render
    has_table = any(b["type"] == "table" for b in blocks)
    has_list = any(b["type"] in ("bullet_list", "numbered_list")
                   for b in blocks)

    if has_table:
        # Prioritize table rendering
        _render_table_on_slide(slide, blocks, content_left, content_top,
                               content_width, content_height)
    else:
        # Text content with bullet points
        _render_text_on_slide(slide, blocks, content_left, content_top,
                              content_width, content_height)


def _render_text_on_slide(slide, blocks: List[Dict],
                          left, top, width, max_height) -> None:
    """Render text blocks (paragraphs, lists, sub-headings) on a slide."""
    TEXT_DARK = PptxRGBColor(0x33, 0x33, 0x33)
    ACCENT = PptxRGBColor(0x2E, 0x75, 0xB6)

    txBox = slide.shapes.add_textbox(left, top, width, max_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first_para = True
    for block in blocks:
        btype = block["type"]

        if btype == "paragraph":
            text = _strip_inline_md(block["text"])
            if not text.strip():
                continue
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            p.text = text
            p.space_after = PptxPt(8)
            p.space_before = PptxPt(4)
            _pptx_set_font(p.add_run(), FONT_BODY_CN, 16, False, TEXT_DARK)

        elif btype == "heading":
            # Sub-heading (### or ####)
            text = _strip_inline_md(block["text"])
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            p.text = text
            p.space_before = PptxPt(14)
            p.space_after = PptxPt(6)
            _pptx_set_font(p.add_run(), FONT_HEADING_CN, 20, True, ACCENT)

        elif btype == "bullet_list":
            for item in block["items"]:
                text = _strip_inline_md(item)
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                p.text = f"  {text}"
                p.level = 0
                p.space_after = PptxPt(4)
                # Bullet character
                pPr = p._pPr
                if pPr is None:
                    from pptx.oxml.ns import qn as pptx_qn
                    pPr = p._p.get_or_add_pPr()
                _pptx_set_font(p.add_run(), FONT_BODY_CN, 15, False, TEXT_DARK)

        elif btype == "numbered_list":
            for idx, item in enumerate(block["items"], 1):
                text = _strip_inline_md(item)
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                p.text = f"  {idx}. {text}"
                p.space_after = PptxPt(4)
                _pptx_set_font(p.add_run(), FONT_BODY_CN, 15, False, TEXT_DARK)

        elif btype == "code_block":
            text = block["code"]
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            p.text = text
            p.space_before = PptxPt(8)
            p.space_after = PptxPt(8)
            _pptx_set_font(p.add_run(), FONT_CODE, 11, False,
                           PptxRGBColor(0x33, 0x33, 0x33))


def _render_table_on_slide(slide, blocks: List[Dict],
                           left, top, width, max_height) -> None:
    """Render the first table found in blocks on the slide."""
    TEXT_DARK = PptxRGBColor(0x33, 0x33, 0x33)
    TEXT_WHITE = PptxRGBColor(0xFF, 0xFF, 0xFF)
    DARK_BG = PptxRGBColor(0x1F, 0x3A, 0x5F)

    # Find first table
    table_block = None
    for b in blocks:
        if b["type"] == "table":
            table_block = b
            break

    if not table_block:
        return

    headers = table_block["headers"]
    rows = table_block["rows"]
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    # Limit rows to fit on slide
    max_slide_rows = min(num_rows, 12)
    if num_rows > max_slide_rows:
        rows = rows[:max_slide_rows - 1]

    # Calculate table dimensions
    table_width = min(width, PptxInches(12.0))
    row_height = PptxInches(0.4)
    table_height = row_height * (len(rows) + 1)

    # Create table shape
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, num_cols, left, top, table_width, table_height)
    table = tbl_shape.table

    # Style header row
    for j, header_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = _strip_inline_md(header_text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                _pptx_set_font(run, FONT_HEADING_CN, 12, True, TEXT_WHITE)

    # Style data rows
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            if j >= num_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = _strip_inline_md(cell_text)
            # Alternating row color
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptxRGBColor(0xED, 0xF2, 0xF9)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _pptx_set_font(run, FONT_BODY_CN, 11, False, TEXT_DARK)

    # Add text content below table if there are non-table blocks
    non_table_blocks = [b for b in blocks if b["type"] != "table"]
    if non_table_blocks:
        text_top = top + table_height + PptxInches(0.2)
        remaining_height = max_height - table_height - PptxInches(0.2)
        if remaining_height > PptxInches(0.5):
            _render_text_on_slide(slide, non_table_blocks,
                                  left, text_top,
                                  width, remaining_height)


# ===========================================================================
# File Routing Logic
# ===========================================================================

def determine_output_format(md_path: Path) -> Optional[str]:
    """
    Determine the output format for a given markdown file path.
    Returns 'docx', 'xlsx', 'pptx', or None (skip).
    """
    filename = md_path.name
    rel_path = str(md_path.relative_to(BASE_DIR)).replace("\\", "/")

    # Skip protected files
    if filename in SKIP_FILES:
        return None

    # PPTX takes priority for specific files
    if rel_path in PPTX_FILES:
        return "pptx"

    # Check XLSX directories
    for xlsx_dir in XLSX_DIRS:
        if rel_path.startswith(xlsx_dir + "/"):
            return "xlsx"

    # Default to DOCX for any other markdown file
    return "docx"


# ===========================================================================
# Main Execution
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Convert Lean knowledge base Markdown files to Office formats.")
    parser.add_argument("--dry-run", action="store_true",
                        help="Show what would be converted without actually converting.")
    parser.add_argument("--keep-md", action="store_true",
                        help="Keep original .md files after conversion.")
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Show detailed progress information.")
    args = parser.parse_args()

    print("=" * 70)
    print("  Lean Knowledge Base: Markdown to Office Converter")
    print("=" * 70)
    print(f"  Base directory: {BASE_DIR}")
    print(f"  Dry run:        {args.dry_run}")
    print(f"  Keep .md:       {args.keep_md}")
    print()

    # Collect all markdown files
    md_files = sorted(BASE_DIR.rglob("*.md"))
    print(f"Found {len(md_files)} markdown files total.")

    # Categorize files
    docx_files: List[Path] = []
    xlsx_files: List[Path] = []
    pptx_files: List[Path] = []
    skipped_files: List[Path] = []

    for md_path in md_files:
        fmt = determine_output_format(md_path)
        if fmt is None:
            skipped_files.append(md_path)
        elif fmt == "docx":
            docx_files.append(md_path)
        elif fmt == "xlsx":
            xlsx_files.append(md_path)
        elif fmt == "pptx":
            pptx_files.append(md_path)

    print(f"  -> Word (.docx): {len(docx_files)} files")
    print(f"  -> Excel (.xlsx): {len(xlsx_files)} files")
    print(f"  -> PPTX (.pptx):  {len(pptx_files)} files")
    print(f"  -> Skipped:       {len(skipped_files)} files")
    print()

    if args.dry_run:
        print("--- DRY RUN: Files that would be converted ---")
        print()

        if docx_files:
            print("Word (.docx) files:")
            for f in docx_files:
                print(f"  {f.relative_to(BASE_DIR)}")
            print()

        if xlsx_files:
            print("Excel (.xlsx) files:")
            for f in xlsx_files:
                print(f"  {f.relative_to(BASE_DIR)}")
            print()

        if pptx_files:
            print("PowerPoint (.pptx) files:")
            for f in pptx_files:
                print(f"  {f.relative_to(BASE_DIR)}")
            print()

        if skipped_files:
            print("Skipped files:")
            for f in skipped_files:
                print(f"  {f.relative_to(BASE_DIR)}")
            print()

        print("Dry run complete. No files were modified.")
        return

    # -- Perform conversions --
    success_count = 0
    error_count = 0
    errors: List[Tuple[Path, str]] = []

    total = len(docx_files) + len(xlsx_files) + len(pptx_files)
    current = 0

    # Convert to Word
    for md_path in docx_files:
        current += 1
        output_path = md_path.with_suffix(".docx")
        progress = f"[{current}/{total}]"
        try:
            print(f"{progress} Converting (DOCX): "
                  f"{md_path.relative_to(BASE_DIR)} ... ", end="", flush=True)
            convert_to_docx(md_path, output_path)
            print("OK")
            success_count += 1
            if not args.keep_md:
                md_path.unlink()
                if args.verbose:
                    print(f"         Deleted: {md_path.name}")
        except Exception as e:
            print(f"FAILED: {e}")
            error_count += 1
            errors.append((md_path, str(e)))

    # Convert to Excel
    for md_path in xlsx_files:
        current += 1
        output_path = md_path.with_suffix(".xlsx")
        progress = f"[{current}/{total}]"
        try:
            print(f"{progress} Converting (XLSX): "
                  f"{md_path.relative_to(BASE_DIR)} ... ", end="", flush=True)
            convert_to_xlsx(md_path, output_path)
            print("OK")
            success_count += 1
            if not args.keep_md:
                md_path.unlink()
                if args.verbose:
                    print(f"         Deleted: {md_path.name}")
        except Exception as e:
            print(f"FAILED: {e}")
            error_count += 1
            errors.append((md_path, str(e)))

    # Convert to PPTX
    for md_path in pptx_files:
        current += 1
        output_path = md_path.with_suffix(".pptx")
        progress = f"[{current}/{total}]"
        try:
            print(f"{progress} Converting (PPTX): "
                  f"{md_path.relative_to(BASE_DIR)} ... ", end="", flush=True)
            convert_to_pptx(md_path, output_path)
            print("OK")
            success_count += 1
            if not args.keep_md:
                md_path.unlink()
                if args.verbose:
                    print(f"         Deleted: {md_path.name}")
        except Exception as e:
            print(f"FAILED: {e}")
            error_count += 1
            errors.append((md_path, str(e)))

    # -- Summary --
    print()
    print("=" * 70)
    print("  Conversion Summary")
    print("=" * 70)
    print(f"  Total files processed: {total}")
    print(f"  Successful:            {success_count}")
    print(f"  Failed:                {error_count}")
    print(f"  Skipped (README/CLAUDE): {len(skipped_files)}")

    if errors:
        print()
        print("  Errors:")
        for path, err in errors:
            print(f"    - {path.relative_to(BASE_DIR)}: {err}")

    print()
    print("Done.")


if __name__ == "__main__":
    main()
