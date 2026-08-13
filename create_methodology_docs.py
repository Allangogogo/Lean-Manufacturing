#!/usr/bin/env python3
"""
Lean Knowledge Base - Methodology Documents Generator
Creates 6 Word documents (.docx) and 2 Training PPTs (.pptx)
for the World-Class Manufacturing Methodology section.
"""

import os
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Path setup
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).resolve().parent
DOCS_DIR = BASE_DIR / "01-精益工具知识库" / "05-世界级制造方法论"
PPT_DIR = BASE_DIR / "02-精益培训" / "02-培训材料"

# ---------------------------------------------------------------------------
# Font constants
# ---------------------------------------------------------------------------
FONT_BODY = "SimSun"       # 宋体
FONT_HEADING = "SimHei"    # 黑体
FONT_BODY_SIZE = 11
FONT_H1_SIZE = 16
FONT_H2_SIZE = 14
FONT_H3_SIZE = 12

# ---------------------------------------------------------------------------
# python-docx helpers
# ---------------------------------------------------------------------------
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn


def _set_cell_shading(cell, color_hex: str):
    """Set cell background shading."""
    shading = cell._element.get_or_add_tcPr()
    shd = shading.makeelement(qn("w:shd"), {
        qn("w:val"): "clear",
        qn("w:color"): "auto",
        qn("w:fill"): color_hex,
    })
    shading.append(shd)


def _set_run_font(run, font_name=FONT_BODY, size=FONT_BODY_SIZE, bold=False, color=None):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    if color:
        run.font.color.rgb = RGBColor(*color)


def add_title(doc: Document, text: str):
    p = doc.add_heading(level=0)
    run = p.add_run(text)
    _set_run_font(run, FONT_HEADING, 22, bold=True)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER


def add_heading1(doc: Document, text: str):
    p = doc.add_heading(level=1)
    run = p.add_run(text)
    _set_run_font(run, FONT_HEADING, FONT_H1_SIZE, bold=True)


def add_heading2(doc: Document, text: str):
    p = doc.add_heading(level=2)
    run = p.add_run(text)
    _set_run_font(run, FONT_HEADING, FONT_H2_SIZE, bold=True)


def add_heading3(doc: Document, text: str):
    p = doc.add_heading(level=3)
    run = p.add_run(text)
    _set_run_font(run, FONT_HEADING, FONT_H3_SIZE, bold=True)


def add_paragraph(doc: Document, text: str, bold=False, indent=False):
    p = doc.add_paragraph()
    run = p.add_run(text)
    _set_run_font(run, FONT_BODY, FONT_BODY_SIZE, bold=bold)
    if indent:
        p.paragraph_format.left_indent = Cm(1)
    return p


def add_bullet(doc: Document, text: str, level: int = 0):
    p = doc.add_paragraph(style="List Bullet")
    p.clear()
    run = p.add_run(text)
    _set_run_font(run, FONT_BODY, FONT_BODY_SIZE)
    if level > 0:
        p.paragraph_format.left_indent = Cm(1.5 * level)
    return p


def add_table(doc: Document, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    # Header row
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ""
        run = cell.paragraphs[0].add_run(h)
        _set_run_font(run, FONT_HEADING, FONT_BODY_SIZE, bold=True, color=(255, 255, 255))
        _set_cell_shading(cell, "1e3a5f")
    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.rows[r_idx + 1].cells[c_idx]
            cell.text = ""
            run = cell.paragraphs[0].add_run(str(val))
            _set_run_font(run, FONT_BODY, FONT_BODY_SIZE - 1)
            if r_idx % 2 == 1:
                _set_cell_shading(cell, "e8eef5")
    doc.add_paragraph()  # spacing
    return table


def add_page_break(doc: Document):
    doc.add_page_break()


# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DARK_BLUE = PRGBColor(0x1e, 0x3a, 0x5f)
DARK_GREEN = PRGBColor(0x06, 0x5f, 0x46)
WHITE = PRGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = PRGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = PRGBColor(0x33, 0x33, 0x33)
MID_BLUE = PRGBColor(0x2c, 0x5f, 0x8a)
ACCENT_ORANGE = PRGBColor(0xE8, 0x6C, 0x00)
TABLE_HEADER_BG = PRGBColor(0x1e, 0x3a, 0x5f)
TABLE_ALT_BG = PRGBColor(0xE8, 0xEE, 0xF5)


def _ppt_set_font(run, name=FONT_HEADING, size=18, bold=False, color=DARK_TEXT):
    from lxml import etree
    run.font.name = name
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Set East Asian font via XML (use _r instead of _element for pptx runs)
    rPr = run._r.get_or_add_rPr()
    rFonts_attr = qn("a:rFonts")
    rFonts_elem = rPr.find(rFonts_attr)
    if rFonts_elem is None:
        rFonts_elem = etree.SubElement(rPr, rFonts_attr)
    rFonts_elem.set(qn("a:latin"), name)
    rFonts_elem.set(qn("a:ea"), name)


def ppt_add_header_bar(slide, color=DARK_BLUE, height=1.1):
    """Add a colored header bar at the top of a slide."""
    from pptx.util import Inches as PInches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        PInches(0), PInches(0), PInches(13.33), PInches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def ppt_add_title_text(slide, text, left=0.6, top=0.15, width=12, height=0.8,
                       font_size=30, color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top),
                                     PInches(width), PInches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _ppt_set_font(run, FONT_HEADING, font_size, bold, color)
    return txBox


def ppt_add_content_box(slide, lines, left=0.6, top=1.4, width=12, height=5.5,
                        font_size=16, bullet=True):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top),
                                     PInches(width), PInches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _ppt_set_font(run, FONT_BODY, font_size, bold=False, color=DARK_TEXT)
        p.space_after = PPt(6)
        if bullet and line.startswith("- "):
            p.level = 0
    return txBox


def ppt_add_table_slide(slide, headers, rows, left=0.5, top=1.5, width=12.3, row_h=0.42):
    cols = len(headers)
    total_rows = 1 + len(rows)
    tbl_shape = slide.shapes.add_table(
        total_rows, cols, PInches(left), PInches(top),
        PInches(width), PInches(row_h * total_rows)
    )
    tbl = tbl_shape.table
    # Header
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        _ppt_set_font(run, FONT_HEADING, 13, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            _ppt_set_font(run, FONT_BODY, 11, bold=False, color=DARK_TEXT)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
    return tbl_shape


# ===================================================================
# DOCUMENT 1: TPS Toyota Production System
# ===================================================================
def create_tps_doc():
    doc = Document()
    add_title(doc, "TPS 丰田生产方式详解")
    add_paragraph(doc, "Toyota Production System -- 世界级制造的基石")
    add_paragraph(doc, "文档编号: WCM-METH-001  |  版本: 1.0  |  分类: 方法论详解")
    add_page_break(doc)

    # --- 1. TPS 定义与历史 ---
    add_heading1(doc, "1. TPS 定义与历史")
    add_heading2(doc, "1.1 什么是 TPS")
    add_paragraph(doc,
        "TPS（Toyota Production System，丰田生产方式）是丰田汽车公司在长期实践中发展出的一套完整的"
        "生产管理哲学和方法体系。它以\"消除浪费\"为核心理念，通过持续改善和尊重人性，追求最优的"
        "生产效率和产品质量。TPS 被公认为现代精益生产（Lean Manufacturing）的直接源头。")
    add_paragraph(doc,
        "TPS 的本质不仅仅是工具和方法的集合，更是一种企业文化和思维方式。它强调全员参与、"
        "现场主义（Genchi Genbutsu）和持续改善（Kaizen），构建了一个能够自我进化的生产系统。")

    add_heading2(doc, "1.2 TPS 的历史演进")
    add_heading3(doc, "创始人：丰田佐吉（1867-1930）")
    add_paragraph(doc,
        "丰田佐吉是丰田集团的创始人，被誉为\"日本发明之父\"。他发明了自动织布机，"
        "并在1924年成功研制出全球首台全自动织布机（G型自动织布机）。他的核心创新在于"
        "引入了\"自働化\"（Jidoka）概念——机器在检测到异常时能够自动停止，从而防止不良品"
        "流入下一道工序。这一理念后来成为 TPS 的两大支柱之一。")
    add_paragraph(doc,
        "佐吉的三大发明精神：① 不依赖他人的独立思考；② 基于现场的实践创新；"
        "③ 将人的智慧赋予机器。这三点奠定了 TPS 的哲学基础。")

    add_heading3(doc, "集大成者：大野耐一（1912-1990）")
    add_paragraph(doc,
        "大野耐一被誉为\"TPS之父\"，他在丰田汽车公司系统性地发展了 TPS。"
        "他于1950年代开始在丰田元町工厂推行看板系统和准时化生产，创造性地将丰田佐吉的"
        "自働化理念与 JIT（准时化）结合，形成了完整的 TPS 体系。")
    add_paragraph(doc, "大野耐一的核心贡献包括：")
    add_bullet(doc, "发明看板（Kanban）系统——实现拉动式生产的信息工具")
    add_bullet(doc, "创立\"大野耐一圆圈\"训练法——培养问题意识和观察力")
    add_bullet(doc, "推行\"后工序拉动\"——从推动式到拉动式的根本转变")
    add_bullet(doc, "倡导\"问5个为什么\"——追溯根本原因的问题分析法")
    add_bullet(doc, "建立「一个流」（One-Piece Flow）生产模式")

    add_heading3(doc, "推广者：藤本三郎（1935-）")
    add_paragraph(doc,
        "藤本三郎是将 TPS 系统化并向全球推广的关键人物。他在1970-80年代担任丰田生产调查部"
        "负责人，将 TPS 的隐性知识转化为显性的管理原则和方法论。他主导了对供应商和海外工厂的"
        " TPS 导入工作，是 TPS 向全球传播的重要推手。")

    add_heading2(doc, "1.3 TPS 的全球影响")
    add_paragraph(doc,
        "1990年，James Womack 等人出版《改变世界的机器》（The Machine That Changed the World），"
        "首次系统介绍了丰田生产方式，并将其命名为\"精益生产\"（Lean Production）。此后，TPS 的"
        "理念迅速传播到全球制造业和服务业，成为当今最具影响力的生产管理方法论之一。")

    add_page_break(doc)

    # --- 2. TPS 屋型图 ---
    add_heading1(doc, "2. TPS 屋型图（TPS House）")
    add_paragraph(doc,
        "TPS 屋型图是理解 TPS 结构最经典的可视化工具。它将 TPS 比作一座建筑，"
        "清晰展示了各要素之间的支撑关系。")

    add_heading2(doc, "2.1 屋型图结构")
    add_paragraph(doc, "【屋顶】-- 最终目标", bold=True)
    add_bullet(doc, "最高质量（Highest Quality）")
    add_bullet(doc, "最低成本（Lowest Cost）")
    add_bullet(doc, "最短交期（Shortest Lead Time）")
    add_bullet(doc, "最佳安全性（Best Safety）")
    add_bullet(doc, "最高士气（Highest Morale）")

    add_paragraph(doc, "【左支柱】-- 自働化（Jidoka）", bold=True)
    add_bullet(doc, "人机分离：机器自动检测异常并停止")
    add_bullet(doc, "安灯系统（Andon）：可视化异常信号")
    add_bullet(doc, "防错（Poka-Yoke）：防止人为失误的装置")
    add_bullet(doc, "异常停止线（Jidoka Line）：确保质量内建于过程")

    add_paragraph(doc, "【右支柱】-- 准时化（JIT）", bold=True)
    add_bullet(doc, "节拍时间（Takt Time）：匹配客户需求的节奏")
    add_bullet(doc, "连续流（Continuous Flow）：消除等待和库存")
    add_bullet(doc, "拉动系统（Pull System）：按需生产")
    add_bullet(doc, "平准化（Heijunka）：均衡化生产排程")

    add_paragraph(doc, "【地基】-- 稳定化与标准化", bold=True)
    add_bullet(doc, "标准作业（Standardized Work）")
    add_bullet(doc, "5S 现场管理")
    add_bullet(doc, "目视管理（Visual Management）")
    add_bullet(doc, "全员生产维护（TPM）")
    add_bullet(doc, "持续改善（Kaizen）")
    add_bullet(doc, "人才培养（Human Development）")

    add_page_break(doc)

    # --- 3. 两大支柱详解 ---
    add_heading1(doc, "3. 两大支柱详解")

    add_heading2(doc, "3.1 自働化（Jidoka）")
    add_paragraph(doc,
        "\"自働化\"中的\"働\"字特指加入了\"人\"字旁的动，意为\"赋予人智慧的自动化\"。"
        "它不仅仅是机器自动运行，更强调在自动化过程中融入人的判断力和智慧。")

    add_heading3(doc, "3.1.1 人机分离")
    add_paragraph(doc,
        "人机分离是自働化的核心原则。传统生产中，操作员必须全程监控机器运行，造成人力浪费。"
        "自働化要求机器具备自动检测异常并停止的能力，使一个操作员能够同时管理多台设备。"
        "这既提高了人员利用率，又确保了产品质量。")

    add_heading3(doc, "3.1.2 异常停止（Abnormal Stop）")
    add_paragraph(doc,
        "当设备检测到异常（如缺料、尺寸偏差、设备故障）时，自动停止运行。"
        "异常停止不是故障停机，而是有意识的、预防性的停机行为。停机后立即发出信号，"
        "等待问题解决后方可恢复生产。这确保了不合格品不会传递到下道工序。")

    add_heading3(doc, "3.1.3 防错（Poka-Yoke）")
    add_paragraph(doc,
        "防错是由新乡重夫提出的概念，目的是通过设计装置或方法，从根本上防止人为失误的发生。"
        "防错分为三类：")
    add_bullet(doc, "预防型防错（Prevention Poka-Yoke）：使错误不可能发生")
    add_bullet(doc, "检测型防错（Detection Poka-Yoke）：使错误在发生时立即被发现")
    add_bullet(doc, "警告型防错（Warning Poka-Yoke）：通过声光信号提醒操作者")
    add_paragraph(doc, "紧固件行业应用：螺纹通止规、扭矩限制器、颜色编码料盒、光电传感器检测漏装等。")

    add_heading3(doc, "3.1.4 安灯系统（Andon）")
    add_paragraph(doc,
        "安灯系统是一种可视化管理工具，通过信号灯（绿-黄-红）显示生产线状态。"
        "当操作员发现问题时，拉动安灯绳或按下按钮，黄色灯亮起表示需要支援，"
        "红色灯亮起表示生产线停止。班组长在规定时间内响应，解决不了则逐级上报。"
        "安灯系统的关键在于赋予一线员工停止生产线的权力，体现了对质量的绝对承诺。")

    add_page_break(doc)

    add_heading2(doc, "3.2 准时化（JIT - Just In Time）")
    add_paragraph(doc,
        "准时化是指在需要的时间、将需要的数量、生产需要的产品。其核心是消除一切浪费，"
        "特别是过量生产和库存浪费。JIT 是 TPS 中最具挑战性的支柱，因为它要求整个生产系统"
        "高度协调和稳定。")

    add_heading3(doc, "3.2.1 节拍时间（Takt Time）")
    add_paragraph(doc,
        "节拍时间是根据客户需求计算的生产节奏，公式为：")
    add_paragraph(doc, "节拍时间 = 可用工作时间 / 客户需求量", bold=True, indent=True)
    add_paragraph(doc,
        "例如：每天可用工作时间 480 分钟，客户需求 240 件，则节拍时间 = 2 分钟/件。"
        "这意味着每 2 分钟必须产出一件产品。节拍时间是整个生产系统的\"心跳\"，"
        "所有工序必须围绕节拍时间来设计和优化。")

    add_heading3(doc, "3.2.2 连续流（Continuous Flow）")
    add_paragraph(doc,
        "连续流是指产品在各工序之间无等待、无中断地流动。实现连续流的关键是：")
    add_bullet(doc, "缩短换型时间（SMED）：使小批量生产经济可行")
    add_bullet(doc, "单元化生产（Cell Manufacturing）：将设备按工艺流程排列")
    add_bullet(doc, "消除瓶颈：平衡各工序的节拍时间")
    add_bullet(doc, "减少在制品库存：暴露隐藏的生产问题")

    add_heading3(doc, "3.2.3 拉动系统（Pull System）")
    add_paragraph(doc,
        "拉动系统是通过看板（Kanban）等工具，由后工序向前工序发出生产指令。"
        "只有当下游工序消耗了物料，上游工序才开始补充生产。这与传统的推动式（Push）生产"
        "形成鲜明对比：推动式根据预测生产，容易造成过量生产和库存积压。"
        "拉动系统确保只生产被需要的产品，从根本上消除过量生产的浪费。")

    add_heading3(doc, "3.2.4 平准化（Heijunka）")
    add_paragraph(doc,
        "平准化是将生产任务在时间和数量上均匀分配，避免大批量集中生产造成的波动。"
        "平准化板（Heijunka Box）是实现平准化的常用工具，通过在时间槽中分配生产看板，"
        "将不同产品的生产均匀地散布在一天的工作时间内。平准化不仅适用于生产排程，"
        "还延伸到物料供应和人员配置的均衡化。")

    add_page_break(doc)

    # --- 4. TPS 14项管理原则 ---
    add_heading1(doc, "4. TPS 14 项管理原则")
    add_paragraph(doc,
        "Jeffrey Liker 在《丰田模式》（The Toyota Way）中将 TPS 总结为 14 项管理原则，"
        "分为 4 大类：")

    add_heading2(doc, "4.1 长期理念（Long-Term Philosophy）")
    add_table(doc,
        ["原则", "内容", "说明"],
        [
            ["原则1", "基于长期理念做决策", "即使牺牲短期财务目标，也要坚持长期发展。"
             "做正确的事，而非仅追求短期利润。"],
        ])

    add_heading2(doc, "4.2 正确的流程产生正确的结果（The Right Process Will Produce the Right Results）")
    add_table(doc,
        ["原则", "内容", "说明"],
        [
            ["原则2", "创造连续流程以暴露问题", "用连续流替代批量处理，使问题无处藏身"],
            ["原则3", "使用拉动系统避免过量生产", "通过看板和拉动减少库存和等待"],
            ["原则4", "均衡化工作负荷（平准化）", "使工作节奏稳定可预测"],
            ["原则5", "建立立即暂停以解决问题的文化（自働化）", "宁可停线也不放过不良品"],
            ["原则6", "标准化任务是持续改善的基石", "先稳定，再改善"],
            ["原则7", "使用可视化管理使问题无处藏身", "一目了然的状态显示"],
            ["原则8", "只使用可靠的、经过验证的技术", "技术服务于人，而非替代人"],
        ])

    add_heading2(doc, "4.3 培养人才与合作伙伴（Develop Your People and Partners）")
    add_table(doc,
        ["原则", "内容", "说明"],
        [
            ["原则9", "培养深度理解公司理念的领导者", "领导者要亲身实践，而非仅靠指令"],
            ["原则10", "培养优秀的人才和团队", "尊重个人发展和团队协作"],
            ["原则11", "尊重合作伙伴和供应商", "将合作伙伴视为延伸，共同成长"],
        ])

    add_heading2(doc, "4.4 持续解决问题以驱动学习（Continuously Solving Problems Drives Learning）")
    add_table(doc,
        ["原则", "内容", "说明"],
        [
            ["原则12", "亲临现场以彻底了解情况（现地现物）", "到现场看实物，了解实际情况"],
            ["原则13", "共识决策，快速执行（根子协商）", "充分讨论但不拖延决策"],
            ["原则14", "通过改善和学习成为学习型组织", "持续改进，永不停步"],
        ])

    add_page_break(doc)

    # --- 5. TPS 与精益制造 ---
    add_heading1(doc, "5. TPS 与精益制造的关系")
    add_paragraph(doc,
        "精益制造（Lean Manufacturing）是在 TPS 基础上发展而来的全球性生产管理理念。"
        "两者的关系可以概括为：")
    add_bullet(doc, "TPS 是精益制造的\"源头\"和\"原型\"")
    add_bullet(doc, "精益制造是 TPS 在全球范围内的\"翻译\"和\"推广\"")
    add_bullet(doc, "精益制造在 TPS 基础上增加了更多工具和方法（如价值流图 VSM）")
    add_bullet(doc, "两者核心理念一致：消除浪费、持续改善、尊重人")

    add_paragraph(doc, "主要区别：", bold=True)
    add_table(doc,
        ["维度", "TPS", "精益制造"],
        [
            ["起源", "丰田汽车公司内部实践", "MIT 研究团队全球推广"],
            ["范围", "主要针对制造业", "扩展到服务业、医疗、软件等"],
            ["工具集", "自働化 + JIT + 看板", "增加 VSM、A3、价值分析等"],
            ["文化基础", "日本企业文化（集体主义）", "需要本地化适应"],
            ["关注点", "系统整体优化", "价值流端到端优化"],
        ])

    add_page_break(doc)

    # --- 6. TPS 在紧固件行业的应用 ---
    add_heading1(doc, "6. TPS 在紧固件行业的应用框架")
    add_paragraph(doc,
        "紧固件行业（螺栓、螺母、垫圈、螺丝等）具有多品种、大批量、连续制造的特点，"
        "TPS 的许多原则可以直接应用。以下是紧固件行业应用 TPS 的框架：")

    add_heading2(doc, "6.1 自働化在紧固件行业的应用")
    add_bullet(doc, "在线检测：利用机器视觉和激光测量技术，实时检测螺纹尺寸、头部形状、表面缺陷")
    add_bullet(doc, "异常停止：当检测到尺寸超差时自动停机，防止批量不良")
    add_bullet(doc, "防错装置：使用模具编号标识、料盒防错、换型确认清单")
    add_bullet(doc, "安灯系统：在冷镦、搓丝、热处理等关键工序设置安灯呼叫点")

    add_heading2(doc, "6.2 JIT 在紧固件行业的应用")
    add_bullet(doc, "节拍时间：根据客户订单计算各工序的生产节拍")
    add_bullet(doc, "连续流：将冷镦→搓丝→热处理→表面处理→检验包装尽量安排在连续流程中")
    add_bullet(doc, "看板拉动：成品仓库向生产线发出看板信号，按实际消耗生产")
    add_bullet(doc, "平准化排程：将不同规格的产品均匀分配到各班次，减少换型次数")

    add_heading2(doc, "6.3 标准化作业在紧固件行业")
    add_bullet(doc, "冷镦工序标准作业：模具安装参数、润滑频次、首件检验流程")
    add_bullet(doc, "热处理工序标准作业：温度曲线、时间控制、冷却介质管理")
    add_bullet(doc, "检验工序标准作业：抽样方案、测量方法、判定标准")

    add_page_break(doc)

    # --- 7. TPS 实施的关键成功因素 ---
    add_heading1(doc, "7. TPS 实施的关键成功因素")

    add_heading2(doc, "7.1 高层领导承诺")
    add_paragraph(doc,
        "TPS 的实施必须获得最高管理层的坚定支持。这不仅仅是口头承诺，更需要领导者"
        "亲身参与现场改善活动，投入资源，并坚持长期推进。丰田的成功经验表明，"
        "没有高层的持续关注，TPS 很容易沦为形式主义。")

    add_heading2(doc, "7.2 全员参与文化")
    add_paragraph(doc,
        "TPS 强调\"全员参与\"——每位员工都是改善的主体。企业需要建立鼓励提案、容忍失败、"
        "奖励改善的文化氛围。丰田每年收到超过百万条员工改善建议，这源于其深厚的文化基础。")

    add_heading2(doc, "7.3 稳定化优先")
    add_paragraph(doc,
        "在追求卓越之前，必须先实现稳定。稳定化包括：稳定的人员、稳定的设备、"
        "稳定的质量、稳定的流程。只有在稳定的基础上，改善才有意义。")
    add_bullet(doc, "先推行 5S 和标准化作业")
    add_bullet(doc, "再实施 TPM 确保设备稳定")
    add_bullet(doc, "然后引入看板和拉动系统")
    add_bullet(doc, "最后追求连续流和平准化")

    add_heading2(doc, "7.4 持续改善机制")
    add_paragraph(doc,
        "TPS 不是一个\"一次性项目\"，而是一种持续的经营方式。企业需要建立：")
    add_bullet(doc, "日常改善（Daily Kaizen）：班组长每日主持改善活动")
    add_bullet(doc, "课题改善（A3 Problem Solving）：针对重大问题的系统性解决")
    add_bullet(doc, "突破性改善（Kaikaku）：对流程进行根本性重新设计")
    add_bullet(doc, "改善提案制度：全员参与的建议和实施机制")

    add_heading2(doc, "7.5 人才培养体系")
    add_paragraph(doc,
        "人才是 TPS 最重要的资产。企业需要建立系统的人才培养机制：")
    add_bullet(doc, "TWI（Training Within Industry）：一线主管培训")
    add_bullet(doc, "多能工培养：使员工掌握多种技能，增加人员弹性")
    add_bullet(doc, "问题解决能力培养：通过 \"五个为什么\" 和 A3 思维训练")
    add_bullet(doc, "改善导师制度：资深员工指导新人")

    doc.add_paragraph()
    add_paragraph(doc, "文档结束  |  WCM-METH-001  |  版本 1.0", bold=True)

    # Save
    path = DOCS_DIR / "01-TPS丰田生产方式详解.docx"
    doc.save(str(path))
    return path


# ===================================================================
# DOCUMENT 2: WCM World-Class Manufacturing
# ===================================================================
def create_wcm_doc():
    doc = Document()
    add_title(doc, "WCM 世界级制造")
    add_paragraph(doc, "World-Class Manufacturing -- 全面整合的卓越制造体系")
    add_paragraph(doc, "文档编号: WCM-METH-002  |  版本: 1.0  |  分类: 方法论详解")
    add_page_break(doc)

    add_heading1(doc, "1. WCM 定义与愿景")
    add_paragraph(doc,
        "WCM（World-Class Manufacturing，世界级制造）是由 Richard J. Schonberger 在1980年代"
        "提出，后由意大利 Fiat 集团的 Tata/Joshua 等人系统化发展的综合制造管理体系。"
        "WCM 的愿景是通过全面整合制造系统的各个维度，达到世界级的运营水平。")
    add_paragraph(doc,
        "WCM 的核心理念：不仅仅是\"做精益\"或\"做 TPM\"，而是将精益生产、TPM、全面质量管理、"
        "供应链管理、安全管理、环境管理等多个维度整合为一个统一的框架，实现制造系统的全面卓越。")

    add_heading2(doc, "1.1 WCM 的核心标准")
    add_bullet(doc, "零损失（Zero Loss）：追求零缺陷、零停机、零事故、零库存浪费")
    add_bullet(doc, "全员参与（Total Employee Involvement）：每个员工都是改善的主体")
    add_bullet(doc, "全系统整合（Total System Integration）：消除孤岛，系统优化")
    add_bullet(doc, "持续改善（Continuous Improvement）：永不满足现状")
    add_bullet(doc, "客户导向（Customer Focus）：一切活动以创造客户价值为目标")

    add_page_break(doc)

    add_heading1(doc, "2. WCM 九大支柱详解")
    add_paragraph(doc,
        "WCM 体系包含九大支柱（Pillars），覆盖制造系统的全部关键维度。"
        "每个支柱都有独立的目标、方法和指标，同时支柱之间相互关联、相互支撑。")

    # Pillar 1: AM
    add_heading2(doc, "2.1 支柱一：AM 自主维护（Autonomous Maintenance）")
    add_paragraph(doc,
        "自主维护是 TPM（全员生产维护）的第一支柱，其核心是让操作员承担日常设备维护责任。"
        "目标是实现\"我的设备我维护\"，使操作员能够自主进行设备的清扫、点检、润滑和简单维修。")
    add_paragraph(doc, "自主维护 7 步法：", bold=True)
    add_table(doc,
        ["步骤", "名称", "关键活动"],
        [
            ["Step 1", "初始清扫", "彻底清扫设备，发现并标记异常点"],
            ["Step 2", "对策与改善", "对发现的异常进行修复和改善"],
            ["Step 3", "清扫润滑基准", "制定清扫、润滑、紧固的临时基准"],
            ["Step 4", "总点检", "对设备进行全面点检教育和训练"],
            ["Step 5", "自主点检", "操作员独立进行设备点检"],
            ["Step 6", "工作场所标准化", "整合为 5S + 自主维护标准"],
            ["Step 7", "自主管理", "操作员实现设备的自主管理"],
        ])
    add_paragraph(doc, "关键指标：设备综合效率 OEE、MTBF（平均故障间隔时间）、MTTR（平均修复时间）。")

    # Pillar 2: PM
    add_heading2(doc, "2.2 支柱二：PM 计划维护（Planned Maintenance）")
    add_paragraph(doc,
        "计划维护是通过预防性维护（Preventive Maintenance）和预测性维护（Predictive Maintenance）"
        "策略，将事后维修转变为主动维护。目标是实现零故障、零停机。")
    add_bullet(doc, "预防性维护：按时间周期进行设备检查、更换易损件")
    add_bullet(doc, "预测性维护：利用振动分析、油液分析、红外热成像等技术预判故障")
    add_bullet(doc, "维护计划编制：基于设备关键度制定年度/月度维护日历")
    add_bullet(doc, "备件管理：合理管理备件库存，平衡可用性和成本")

    # Pillar 3: QM
    add_heading2(doc, "2.3 支柱三：QM 质量管理（Quality Management）")
    add_paragraph(doc,
        "质量管理支柱旨在实现零缺陷制造。通过统计过程控制（SPC）、防错装置（Poka-Yoke）、"
        "过程能力分析（Cp/Cpk）等工具，将质量控制从终端检验前移到制造过程。")
    add_bullet(doc, "源头质量管理：在产生缺陷的工序内检出和消除不良")
    add_bullet(doc, "过程能力分析：确保 Cp >= 1.33, Cpk >= 1.0")
    add_bullet(doc, "质量功能展开（QFD）：将客户需求转化为设计和过程参数")
    add_bullet(doc, "8D 问题解决：系统性分析和解决质量问题")
    add_bullet(doc, "SPC 控制图：监控过程稳定性和能力")

    # Pillar 4: DM
    add_heading2(doc, "2.4 支柱四：DM 交期管理（Delivery Management）")
    add_paragraph(doc,
        "交期管理支柱关注准时交货率和交期缩短。通过价值流优化、瓶颈管理、"
        "生产计划优化等手段，实现短交期和高准时交货率。")
    add_bullet(doc, "价值流图分析：识别增值和非增值活动")
    add_bullet(doc, "前置时间缩短：从接单到交付的全流程时间压缩")
    add_bullet(doc, "产线平衡：消除瓶颈工序，提高产出率")
    add_bullet(doc, "交期承诺管理：建立可靠的交期评估和承诺机制")

    # Pillar 5: EM
    add_heading2(doc, "2.5 支柱五：EM 环境管理（Environment Management）")
    add_paragraph(doc,
        "环境管理支柱确保制造过程符合环境法规要求，并持续减少环境影响。"
        "与 ISO 14001 环境管理体系对接。")
    add_bullet(doc, "废弃物减量化：从源头减少废弃物产生")
    add_bullet(doc, "能源管理：降低单位产品能耗")
    add_bullet(doc, "有害物质管控：符合 RoHS、REACH 等法规要求")
    add_bullet(doc, "碳排放管理：碳足迹核算与减碳路径规划")

    # Pillar 6: LM
    add_heading2(doc, "2.6 支柱六：LM 精益管理（Lean Management）")
    add_paragraph(doc,
        "精益管理支柱是 WCM 的效率引擎，将精益生产的核心工具和理念整合到日常管理中。")
    add_bullet(doc, "价值流管理：从客户端到供应商端的全价值流优化")
    add_bullet(doc, "5S 与目视管理：打造整洁、有序、可视化的工作场所")
    add_bullet(doc, "标准作业：将最佳实践固化为标准")
    add_bullet(doc, "快速换型（SMED）：缩短设备换型时间")
    add_bullet(doc, "看板与拉动系统：按需生产和物料流动")
    add_bullet(doc, "全员改善（Kaizen）：持续改善的文化和机制")

    # Pillar 7: SM
    add_heading2(doc, "2.7 支柱七：SM 安全管理（Safety Management）")
    add_paragraph(doc,
        "安全管理支柱以\"零事故\"为目标，通过危险源辨识、风险评估、安全标准化等手段，"
        "建立本质安全的工作环境。")
    add_bullet(doc, "危险源辨识与风险评估（HIRA）")
    add_bullet(doc, "作业安全标准化")
    add_bullet(doc, "安全行为观察（BBS）")
    add_bullet(doc, "应急响应预案")
    add_bullet(doc, "安全培训与意识提升")

    # Pillar 8: OE
    add_heading2(doc, "2.8 支柱八：OE 组织效率（Organization Effectiveness）")
    add_paragraph(doc,
        "组织效率支柱关注人力资源的优化配置和组织能力的提升。")
    add_bullet(doc, "多能工培养：扩大员工技能矩阵")
    add_bullet(doc, "团队管理：班组建设和团队改善活动")
    add_bullet(doc, "绩效管理体系：目标分解与绩效追踪")
    add_bullet(doc, "知识管理：隐性知识显性化，最佳实践共享")

    # Pillar 9: CM
    add_heading2(doc, "2.9 支柱九：CM 成本管理（Cost Management）")
    add_paragraph(doc,
        "成本管理支柱通过全员参与的成本意识提升和系统化的成本降低活动，"
        "实现制造成本的持续优化。")
    add_bullet(doc, "成本可视化：将成本信息传递到一线")
    add_bullet(doc, "损失地图（Loss Map）：识别和量化各种损失")
    add_bullet(doc, "成本改善课题：围绕损失开展专项改善")
    add_bullet(doc, "价值工程（VA/VE）：从设计端降低成本")
    add_bullet(doc, "采购成本管理：供应链协同降本")

    add_page_break(doc)

    add_heading1(doc, "3. WCM 实施阶段")
    add_paragraph(doc, "WCM 实施通常经历 5 个阶段：")

    add_table(doc,
        ["阶段", "名称", "特征", "典型周期"],
        [
            ["阶段1", "导入期", "高层承诺、组织架构建立、培训启动、现状评估", "3-6个月"],
            ["阶段2", "基础期", "5S推行、设备基础维护、标准作业建立", "6-12个月"],
            ["阶段3", "发展期", "各支柱全面推行、改善活动活跃、损失持续降低", "12-24个月"],
            ["阶段4", "成熟期", "系统整合优化、跨部门协同、对标世界级水平", "24-36个月"],
            ["阶段5", "卓越期", "持续自我进化、创新引领、行业标杆", "持续"],
        ])

    add_heading1(doc, "4. WCM 评估模型")
    add_paragraph(doc,
        "WCM 采用分级评估体系，衡量企业在各支柱上的成熟度：")
    add_table(doc,
        ["等级", "全称", "含义", "标准"],
        [
            ["DC", "Development of Compliance", "合规发展", "基础标准已建立，正在推行中"],
            ["JD", "Job Discipline", "工作纪律", "标准执行到位，具备基本纪律性"],
            ["WD", "World-Class Discipline", "世界级纪律", "全面达到世界级标准，持续改善"],
        ])
    add_paragraph(doc,
        "评估方式：WCM 评估通常由内部评估团队或外部专家进行，对每个支柱按5个维度打分："
        "组织与人员、自愿改善活动、信息交换、过程/系统/技术、结果。"
        "每个维度1-5分，总分30分/支柱。")

    add_page_break(doc)

    add_heading1(doc, "5. WCM 与 Lean、TPM 的关系和区别")
    add_table(doc,
        ["维度", "WCM", "Lean", "TPM"],
        [
            ["范围", "全面（9大支柱）", "价值流和流程优化", "设备维护和可靠性"],
            ["起源", "Fiat/Schonberger", "TPS/丰田", "日本维护协会"],
            ["工具集", "Lean + TPM + 质量 + 安全 + 环境 + 成本", "VSM、看板、5S、SMED等", "自主维护、计划维护、OEE等"],
            ["关注点", "全面整合", "消除浪费、加快流动", "设备效率和可靠性"],
            ["评估", "DC/JD/WD分级", "成熟度评估", "三级评价"],
        ])
    add_paragraph(doc,
        "WCM 可以理解为 Lean + TPM + 其他管理维度的超级整合。它以 Lean 和 TPM 为核心，"
        "同时纳入了质量、安全、环境、成本、交期、组织效率等维度，形成完整的制造管理体系。")

    add_page_break(doc)

    add_heading1(doc, "6. WCM 在紧固件行业的应用")
    add_paragraph(doc,
        "紧固件行业的特点——多工序连续制造、设备密集、品种多、质量要求高——"
        "非常适合 WCM 体系的应用。")
    add_bullet(doc, "AM 自主维护：冷镦机、搓丝机、热处理炉的日常点检由操作员负责")
    add_bullet(doc, "PM 计划维护：关键设备（如多工位冷镦机）的预防性维护计划")
    add_bullet(doc, "QM 质量管理：螺纹通止规管理、扭矩强度在线监控、SPC 控制图")
    add_bullet(doc, "DM 交期管理：从冷镦到成品的全流程前置时间压缩")
    add_bullet(doc, "LM 精益管理：5S 推行、快速换模（模具更换时间缩短）、看板系统")
    add_bullet(doc, "SM 安全管理：冲压设备安全防护、化学品管理、噪音防护")
    add_bullet(doc, "CM 成本管理：材料利用率提升、能耗降低、换型损失减少")

    add_page_break(doc)

    add_heading1(doc, "7. WCM 关键指标")
    add_table(doc,
        ["支柱", "关键指标", "计算公式 / 说明"],
        [
            ["AM", "OEE（设备综合效率）", "可用率 x 性能率 x 质量率"],
            ["PM", "MTBF / MTTR", "平均故障间隔 / 平均修复时间"],
            ["QM", "Cpk / 不良率 PPM", "过程能力指数 / 每百万件不良数"],
            ["DM", "准时交货率 OTD", "准时交付订单数 / 总订单数"],
            ["EM", "单位产品能耗", "总能耗 / 产出量"],
            ["LM", "前置时间 Lead Time", "从接单到交付的总时间"],
            ["SM", "事故频率率 / 严重率", "工伤事故次数 / 工作时间"],
            ["OE", "人均产出", "总产出量 / 员工人数"],
            ["CM", "制造成本降低率", "（本期成本 - 基准成本）/ 基准成本"],
        ])

    doc.add_paragraph()
    add_paragraph(doc, "文档结束  |  WCM-METH-002  |  版本 1.0", bold=True)

    path = DOCS_DIR / "02-WCM世界级制造.docx"
    doc.save(str(path))
    return path


# ===================================================================
# DOCUMENT 3: Lean Six Sigma
# ===================================================================
def create_lean_six_sigma_doc():
    doc = Document()
    add_title(doc, "精益六西格玛")
    add_paragraph(doc, "Lean Six Sigma -- 效率与质量的完美结合")
    add_paragraph(doc, "文档编号: WCM-METH-003  |  版本: 1.0  |  分类: 方法论详解")
    add_page_break(doc)

    add_heading1(doc, "1. Lean + Six Sigma 整合方法论")
    add_paragraph(doc,
        "精益六西格玛（Lean Six Sigma, LSS）是将精益生产（Lean）和六西格玛（Six Sigma）"
        "两大方法论有机整合的综合改进方法。Lean 关注速度和流动，消除浪费；Six Sigma 关注"
        "变异和质量，减少缺陷。两者互补结合，实现了效率和质量的双重提升。")

    add_heading2(doc, "1.1 为什么要整合")
    add_paragraph(doc,
        "单独使用 Lean 或 Six Sigma 都有局限性：")
    add_bullet(doc, "Lean 的局限：善于消除浪费和加快流动，但对变异控制和统计分析较弱")
    add_bullet(doc, "Six Sigma 的局限：善于减少变异和提升质量，但对流程速度和浪费关注不够")
    add_bullet(doc, "LSS 的优势：Lean 加速流程 + Six Sigma 减少变异 = 更快、更稳定、更高质量的产出")

    add_heading2(doc, "1.2 Lean 关注速度和流动")
    add_paragraph(doc, "Lean 的核心工具：")
    add_bullet(doc, "价值流图（VSM）：识别流程中的浪费")
    add_bullet(doc, "5S：打造有序的工作环境")
    add_bullet(doc, "看板（Kanban）：实现拉动生产")
    add_bullet(doc, "SMED：快速换型，减少批量")
    add_bullet(doc, "连续流（Continuous Flow）：消除等待和中断")
    add_bullet(doc, "标准作业：固化最佳实践")

    add_heading2(doc, "1.3 Six Sigma 关注变异和质量")
    add_paragraph(doc, "Six Sigma 的核心工具：")
    add_bullet(doc, "统计过程控制（SPC）：监控过程稳定性")
    add_bullet(doc, "过程能力分析（Cp/Cpk/Ppk）：评估过程满足规格的能力")
    add_bullet(doc, "假设检验（t检验、ANOVA等）：验证因果关系")
    add_bullet(doc, "回归分析：建立变量之间的数学关系")
    add_bullet(doc, "实验设计（DOE）：系统优化过程参数")
    add_bullet(doc, "失效模式分析（FMEA）：预防性风险评估")

    add_page_break(doc)

    add_heading1(doc, "2. DMAIC 方法论详解")
    add_paragraph(doc,
        "DMAIC 是精益六西格玛的核心方法论框架，代表五个阶段：Define、Measure、Analyze、"
        "Improve、Control。它为问题解决提供了结构化的路径。")

    add_heading2(doc, "2.1 Define（定义）")
    add_paragraph(doc, "目标：明确问题、目标和范围", bold=True)
    add_bullet(doc, "项目章程（Project Charter）：定义问题陈述、目标、范围、时间表")
    add_bullet(doc, "SIPOC 分析：Supplier-Input-Process-Output-Customer，确定流程边界")
    add_bullet(doc, "客户之声（VOC）：收集和分析客户需求")
    add_bullet(doc, "关键质量特性（CTQ）：将 VOC 转化为可测量的质量指标")
    add_bullet(doc, "高层利害关系者分析：识别项目的支持者和影响者")

    add_heading2(doc, "2.2 Measure（测量）")
    add_paragraph(doc, "目标：收集数据，建立当前过程的基线", bold=True)
    add_bullet(doc, "数据收集计划：确定测量对象、方法、频率、样本量")
    add_bullet(doc, "测量系统分析（MSA/Gage R&R）：验证测量系统的可靠性和重复性")
    add_bullet(doc, "过程能力基线：计算当前的 Cp、Cpk、PPM")
    add_bullet(doc, "价值流图（VSM）：量化流程中的时间、库存和增值比")
    add_bullet(doc, "流程图绘制：详细记录当前流程的每个步骤")

    add_heading2(doc, "2.3 Analyze（分析）")
    add_paragraph(doc, "目标：识别根本原因和关键影响因素", bold=True)
    add_bullet(doc, "鱼骨图（因果图）：从人、机、料、法、环、测六个维度分析原因")
    add_bullet(doc, "5 Why 分析：追溯根本原因")
    add_bullet(doc, "假设检验：用统计方法验证因果假设（t检验、卡方检验、ANOVA）")
    add_bullet(doc, "回归分析：确定 X 和 Y 之间的数学关系")
    add_bullet(doc, "帕累托分析：识别关键少数因素（80/20法则）")
    add_bullet(doc, "FMEA（失效模式与影响分析）：评估潜在风险的严重度、发生度和检出度")

    add_heading2(doc, "2.4 Improve（改进）")
    add_paragraph(doc, "目标：制定和实施改进方案", bold=True)
    add_bullet(doc, "头脑风暴和创意生成：团队协作产生改进方案")
    add_bullet(doc, "实验设计（DOE）：系统地测试和优化关键参数")
    add_bullet(doc, "防错设计（Poka-Yoke）：从设计上防止错误发生")
    add_bullet(doc, "精益工具应用：5S、SMED、看板等")
    add_bullet(doc, "试运行（Pilot）：在小范围内验证改进效果")
    add_bullet(doc, "成本效益分析：评估改进方案的经济可行性")

    add_heading2(doc, "2.5 Control（控制）")
    add_paragraph(doc, "目标：维持改进成果，防止回退", bold=True)
    add_bullet(doc, "控制计划（Control Plan）：定义监控方法、频率和响应措施")
    add_bullet(doc, "SPC 控制图：实时监控过程状态")
    add_bullet(doc, "标准作业文件更新：将改进后的流程固化为标准")
    add_bullet(doc, "培训和交接：确保相关人员掌握新流程")
    add_bullet(doc, "项目成果报告：总结收益和经验教训")
    add_bullet(doc, "过程所有权转移：将监控责任移交给过程负责人")

    add_page_break(doc)

    add_heading1(doc, "3. DFSS 简介")
    add_paragraph(doc,
        "DFSS（Design for Six Sigma，六西格玛设计）是将六西格玛理念应用于新产品和"
        "新流程设计的方法。与 DMAIC（改进现有流程）不同，DFSS 从设计阶段就将质量内建，"
        "目标是在设计阶段就达到六西格玛水平。")
    add_paragraph(doc, "DFSS 常用框架：IDOV（Identify-Design-Optimize-Validate）", bold=True)
    add_bullet(doc, "Identify：识别客户需求和项目范围")
    add_bullet(doc, "Design：设计概念开发和详细设计")
    add_bullet(doc, "Optimize：通过 DOE 等方法优化设计参数")
    add_bullet(doc, "Validate：验证设计是否满足要求")

    add_page_break(doc)

    add_heading1(doc, "4. 西格玛水平计算")
    add_paragraph(doc,
        "西格玛（Sigma）水平衡量过程满足客户规格的能力。1 西格玛水平最低，6 西格玛水平最高。")

    add_table(doc,
        ["西格玛水平", "DPMO（百万机会缺陷数）", "合格率", "含义"],
        [
            ["1 sigma", "691,462", "30.85%", "过程极不稳定"],
            ["2 sigma", "308,538", "69.15%", "过程不稳定"],
            ["3 sigma", "66,807", "93.32%", "一般制造水平"],
            ["4 sigma", "6,210", "99.38%", "较好水平"],
            ["5 sigma", "233", "99.977%", "优秀水平"],
            ["6 sigma", "3.4", "99.99966%", "世界级水平（目标）"],
        ])
    add_paragraph(doc,
        "DPMO = (Defects x 1,000,000) / (Opportunities x Units)")
    add_paragraph(doc,
        "在紧固件行业，典型目标为 4-5 西格玛水平。例如，对于螺纹精度控制，"
        "如果规格为 M10x1.5 的中径公差为 +/-0.05mm，过程标准差为 0.015mm，"
        "则 Cpk = 0.05/(3*0.015) = 1.11，约对应 3.3 西格玛水平。")

    add_page_break(doc)

    add_heading1(doc, "5. 角色体系")
    add_table(doc,
        ["角色", "英文", "职责", "培训时长", "项目要求"],
        [
            ["黄带", "Yellow Belt", "了解 LSS 基础，支持项目团队", "1-2天", "无"],
            ["绿带", "Green Belt", "领导改进项目，运用 DMAIC 工具", "2-3周", "完成1-2个项目"],
            ["黑带", "Black Belt", "全职改进专家，领导复杂项目，培训绿带", "4-6周", "完成4-6个项目"],
            ["大黑带", "Master Black Belt", "战略规划，项目选择，培训黑带/绿带", "持续", "领导组织级 LSS 推行"],
        ])

    add_page_break(doc)

    add_heading1(doc, "6. LSS 在紧固件行业的应用案例")
    add_heading2(doc, "案例：降低六角螺母冷镦工序不良率")
    add_paragraph(doc, "项目背景：", bold=True)
    add_paragraph(doc,
        "某紧固件企业六角螺母冷镦工序月产量 500 万件，不良率 2.8%（PPM = 28,000），"
        "主要缺陷为头部裂纹（45%）和对角尺寸超差（30%）。项目目标将不良率降至 0.5%。")

    add_paragraph(doc, "DMAIC 过程：", bold=True)
    add_bullet(doc, "Define：项目章程定义问题、目标、范围。SIPOC 分析确定过程边界。")
    add_bullet(doc, "Measure：收集 30 天数据，MSA 确认测量系统可靠。当前 Cpk = 0.89。")
    add_bullet(doc, "Analyze：鱼骨图 + 5 Why 找到根本原因——原材料硬度波动大（C=35%）、模具温度控制不稳（C=25%）、润滑不均匀（C=20%）。")
    add_bullet(doc, "Improve：DOE 优化模具温度和润滑参数，建立原材料入库检验标准。")
    add_bullet(doc, "Control：SPC 控制图监控关键参数，标准化作业文件更新，培训操作员。")
    add_paragraph(doc, "项目成果：不良率从 2.8% 降至 0.3%，Cpk 提升至 1.52，年节约成本约 120 万元。")

    add_page_break(doc)

    add_heading1(doc, "7. LSS 与 TPS/WCM 的关系")
    add_table(doc,
        ["维度", "LSS", "TPS", "WCM"],
        [
            ["方法论", "DMAIC 数据驱动", "两大支柱 + 14原则", "9大支柱全面整合"],
            ["改进方式", "项目制、统计驱动", "日常改善 + 课题改善", "支柱推进 + 评估"],
            ["工具偏好", "统计工具为主", "看板、自働化为主", "综合工具"],
            ["人才体系", "GB/BB/MBB", "多能工、TWI", "内部评估员"],
            ["整合建议", "LSS 可作为 WCM/TPS 的问题解决方法论", "", ""],
        ])

    doc.add_paragraph()
    add_paragraph(doc, "文档结束  |  WCM-METH-003  |  版本 1.0", bold=True)

    path = DOCS_DIR / "03-精益六西格玛.docx"
    doc.save(str(path))
    return path


# ===================================================================
# DOCUMENT 4: TOC Theory of Constraints
# ===================================================================
def create_toc_doc():
    doc = Document()
    add_title(doc, "TOC 约束理论")
    add_paragraph(doc, "Theory of Constraints -- 突破瓶颈的系统方法")
    add_paragraph(doc, "文档编号: WCM-METH-004  |  版本: 1.0  |  分类: 方法论详解")
    add_page_break(doc)

    add_heading1(doc, "1. TOC 定义")
    add_paragraph(doc,
        "TOC（Theory of Constraints，约束理论）是由以色列物理学家 Eliyahu M. Goldratt 博士"
        "于1984年在小说《目标》（The Goal）中首次系统阐述的管理理论。TOC 的核心思想是："
        "任何系统在任何给定时刻都存在一个或极少数的约束（瓶颈），系统的产出完全取决于"
        "这些约束。因此，改善系统的关键在于识别和管理系统约束。")
    add_paragraph(doc,
        "TOC 不仅是一种生产管理方法，更是一种全面的管理哲学，应用于生产、项目管理、"
        "供应链、财务、营销等多个领域。")

    add_heading2(doc, "1.1 TOC 的三大核心假设")
    add_bullet(doc, "系统是由相互关联的部分组成的整体")
    add_bullet(doc, "任何系统的产出都受到至少一个约束的限制")
    add_bullet(doc, "所有约束都是可以被打破的——但一次只能打破一个")

    add_page_break(doc)

    add_heading1(doc, "2. 五步聚焦法（Five Focusing Steps）")
    add_paragraph(doc,
        "五步聚焦法是 TOC 的核心方法论，提供了系统性识别和突破约束的路径。")

    add_heading2(doc, "Step 1: 识别约束（Identify the Constraint）")
    add_paragraph(doc,
        "找到系统中限制产出的瓶颈。在制造环境中，瓶颈通常表现为：")
    add_bullet(doc, "堆积最多在制品的工序前")
    add_bullet(doc, "经常加班的工序")
    add_bullet(doc, "交期延迟最频繁的产品对应的工序")
    add_bullet(doc, "产能利用率最高的设备（接近或达到100%）")
    add_paragraph(doc,
        "识别方法：产能负荷分析、在制品分布调查、设备利用率统计、瓶颈工序跟踪。")

    add_heading2(doc, "Step 2: 挖尽约束（Exploit the Constraint）")
    add_paragraph(doc,
        "在不增加投资的前提下，最大化现有约束的产出。确保瓶颈设备100%被有效利用。")
    add_bullet(doc, "消除瓶颈设备的一切停机时间（计划外停机、换型、等待物料）")
    add_bullet(doc, "确保瓶颈设备只加工合格品（不做返工件）")
    add_bullet(doc, "安排最有经验的操作员操作瓶颈设备")
    add_bullet(doc, "利用午休和休息时间让瓶颈设备持续运行")

    add_heading2(doc, "Step 3: 迁就约束（Subordinate Everything Else）")
    add_paragraph(doc,
        "让所有非瓶颈工序配合瓶颈工序的节奏运行。这是TOC最反直觉但最关键的一步。")
    add_bullet(doc, "非瓶颈工序不需要满负荷运转——过度生产只会增加库存")
    add_bullet(doc, "所有生产指令来自瓶颈工序（DBR 系统）")
    add_bullet(doc, "物料投放速度由瓶颈决定，而非预测")

    add_heading2(doc, "Step 4: 打破约束（Elevate the Constraint）")
    add_paragraph(doc,
        "如果挖尽现有能力后约束仍然存在，则投资扩大约束产能。")
    add_bullet(doc, "增加瓶颈设备（购买新设备或增加班次）")
    add_bullet(doc, "外包瓶颈工序")
    add_bullet(doc, "工艺改进提高瓶颈效率")
    add_bullet(doc, "注意：只有在 Step 2-3 用尽后才进行投资")

    add_heading2(doc, "Step 5: 防止惰性（Prevent Inertia）")
    add_paragraph(doc,
        "当约束被打破后，它可能不再是瓶颈。必须回到 Step 1，重新识别新的约束。"
        "不要让旧的习惯和流程成为新的障碍。持续改善的关键在于不断寻找和突破新的约束。")

    add_page_break(doc)

    add_heading1(doc, "3. DBR 排程系统")
    add_paragraph(doc,
        "DBR（Drum-Buffer-Rope）是 TOC 在生产排程中的核心应用。")

    add_heading2(doc, "3.1 Drum（鼓）")
    add_paragraph(doc,
        "鼓是瓶颈工序的节拍。瓶颈工序的生产计划就是整个工厂的\"鼓声\"，"
        "所有其他工序必须跟随这个节拍。瓶颈的排程决定了整个系统的产出。")

    add_heading2(doc, "3.2 Buffer（缓冲）")
    add_paragraph(doc,
        "缓冲是在瓶颈工序前设置的时间缓冲或库存缓冲，确保瓶颈永远不缺料。"
        "缓冲的大小根据系统的波动性确定。如果缓冲经常被耗尽，说明缓冲太小或系统波动太大。")

    add_heading2(doc, "3.3 Rope（绳子）")
    add_paragraph(doc,
        "绳子是连接第一道工序和瓶颈工序的信息通道。它控制原材料的投放时机和数量。"
        "当瓶颈消耗一个单位的物料时，通过绳子信号通知第一道工序投放新的原材料。"
        "这确保了在制品库存始终维持在合理水平。")

    add_page_break(doc)

    add_heading1(doc, "4. 缓冲管理")
    add_paragraph(doc,
        "缓冲管理是 TOC 确保交付可靠性的核心机制。缓冲分为三种类型：")
    add_table(doc,
        ["缓冲类型", "位置", "作用", "管理方法"],
        [
            ["时间缓冲", "瓶颈工序之前", "确保瓶颈不因上游波动而停工", "监控缓冲消耗率，调整大小"],
            ["库存缓冲", "出货前（发运缓冲）", "应对下游需求波动", "设定安全库存水平"],
            ["产能缓冲", "非瓶颈工序", "为约束转移预留弹性", "保持一定的闲置产能"],
        ])
    add_paragraph(doc,
        "缓冲管理的日常操作：每天检查缓冲消耗状态，分为绿色（正常）、黄色（需关注）、"
        "红色（即将耗尽）三个等级，采取相应的行动。")

    add_page_break(doc)

    add_heading1(doc, "5. TOC 三大财务指标")
    add_paragraph(doc,
        "TOC 提出了与传统成本会计不同的财务评价体系：")

    add_table(doc,
        ["指标", "英文", "定义", "目标"],
        [
            [" throughput", "Throughput (T)", "系统通过销售获得收入的速率", "最大化 T"],
            [" inventory", "Inventory (I)", "系统中所有用于销售的资源的总投入", "最小化 I"],
            [" operating expense", "Operating Expense (OE)", "系统将库存转化为产出的花费", "最小化 OE"],
        ])
    add_paragraph(doc,
        "TOC 财务原则：增加 T 比减少 OE 更重要；减少 I 比减少 OE 更容易实现。"
        "传统成本会计追求降低单位成本，可能导致过量生产；TOC 追求最大化整个系统的 Throughput。")

    add_page_break(doc)

    add_heading1(doc, "6. TOC 在紧固件行业的应用")
    add_paragraph(doc,
        "紧固件生产中，热处理工序通常是典型的瓶颈。以下以热处理为约束进行分析：")

    add_heading2(doc, "6.1 识别瓶颈")
    add_paragraph(doc,
        "热处理炉的特点：投资大、升温时间长、批次处理（非连续流）、产能有限。"
        "在紧固件工厂，冷镦和搓丝的速度远高于热处理，导致热处理成为系统瓶颈。")

    add_heading2(doc, "6.2 挖尽热处理产能")
    add_bullet(doc, "减少热处理炉的空炉时间和升降温次数")
    add_bullet(doc, "优化装炉排列，最大化每炉装载量")
    add_bullet(doc, "减少热处理前的等待时间（确保来料及时到达）")
    add_bullet(doc, "安排经验丰富的操作员，减少操作失误导致的返工")
    add_bullet(doc, "通过工艺优化缩短保温时间（在质量允许的范围内）")

    add_heading2(doc, "6.3 迁就约束")
    add_bullet(doc, "冷镦和搓丝工序按热处理的节拍生产，避免过量生产导致库存积压")
    add_bullet(doc, "原材料投放由热处理消耗量拉动")
    add_bullet(doc, "非瓶颈工序在完成热处理需求后可以安排维护或培训")

    add_heading2(doc, "6.4 打破约束")
    add_bullet(doc, "投资新增热处理炉或连续式热处理线")
    add_bullet(doc, "引入真空热处理或感应加热等高效热处理技术")
    add_bullet(doc, "考虑外协热处理（将部分产能外包）")

    add_page_break(doc)

    add_heading1(doc, "7. TOC 与 Lean、TPM 的互补关系")
    add_paragraph(doc,
        "TOC、Lean 和 TPM 并非互相排斥，而是可以互补整合：")

    add_table(doc,
        ["维度", "TOC", "Lean", "TPM"],
        [
            ["核心问题", "瓶颈在哪里？", "浪费在哪里？", "设备可靠吗？"],
            ["改进焦点", "约束点集中改善", "全价值流消除浪费", "设备效率全面提升"],
            ["库存观点", "瓶颈前需要缓冲库存", "库存是浪费，要减少", "库存反映设备问题"],
            ["排程方法", "DBR（鼓-缓冲-绳）", "看板（拉动）", "基于设备能力排程"],
            ["财务视角", "T/I/OE 全局优化", "消除七大浪费", "OEE 提升"],
            ["互补点", "TOC 找到改善重点", "Lean 消除非瓶颈浪费", "TPM 保障设备可靠性"],
        ])
    add_paragraph(doc,
        "整合建议：以 TOC 的五步聚焦法找到系统的约束点（重点改善方向），"
        "在约束点应用 Lean 工具消除浪费、提高流动，同时用 TPM 保障约束设备的可靠性。"
        "对于非约束工序，Lean 的工具（5S、标准作业、防错）同样适用。")

    doc.add_paragraph()
    add_paragraph(doc, "文档结束  |  WCM-METH-004  |  版本 1.0", bold=True)

    path = DOCS_DIR / "04-TOC约束理论.docx"
    doc.save(str(path))
    return path


# ===================================================================
# DOCUMENT 5: Operational Excellence Models
# ===================================================================
def create_excellence_model_doc():
    doc = Document()
    add_title(doc, "卓越运营模型")
    add_paragraph(doc, "Operational Excellence Models -- 从优秀到卓越的框架")
    add_paragraph(doc, "文档编号: WCM-METH-005  |  版本: 1.0  |  分类: 方法论详解")
    add_page_break(doc)

    add_heading1(doc, "1. Shingo 模型（10 大原则）")
    add_paragraph(doc,
        "Shingo Prize（Shingo 研究奖）是由美国犹他州立大学设立的卓越运营大奖，"
        "被誉为\"制造业的诺贝尔奖\"。Shingo 模型基于 Shigeo Shingo（新乡重夫）的"
        "研究成果，提出了 10 大卓越运营原则。")

    add_table(doc,
        ["序号", "原则", "核心含义", "实践要点"],
        [
            ["1", "Lead with humility", "谦逊领导", "领导者到现场倾听、学习，而非仅靠权力指挥"],
            ["2", "Respect every individual", "尊重每个人", "重视每位员工的贡献和潜力"],
            ["3", "Focus on process", "聚焦过程", "结果是过程的产物，改善过程才能改善结果"],
            ["4", "Assure quality at the source", "源头保证质量", "在产生缺陷的地方检出和消除，而非依赖检验"],
            ["5", "Flow and pull value", "流动和拉动价值", "消除中断和等待，让价值连续流动"],
            ["6", "Think systematically", "系统思考", "从全局优化而非局部优化"],
            ["7", "Establish constancy of purpose", "确立恒定目标", "长期一致的方向和承诺"],
            ["8", "Create value for the customer", "为客户创造价值", "一切活动以客户价值为导向"],
            ["9", "Seek perfection", "追求完美", "持续改善，永不满足"],
            ["10", "Think scientifically", "科学思考", "基于数据和实验进行决策和改善"],
        ])

    add_paragraph(doc, "Shingo 模型的关键洞察：", bold=True)
    add_bullet(doc, "理想行为（Ideal Behavior）驱动理想结果（Ideal Results）")
    add_bullet(doc, "系统和工具（System/Tools）支撑理想行为")
    add_bullet(doc, "管理原则（Management Principles）指导系统和工具的设计")
    add_paragraph(doc,
        "Shingo 模型的最深层是管理原则，它决定了企业的文化、系统和最终结果。"
        "仅使用工具而不改变原则，无法实现真正的卓越运营。")

    add_page_break(doc)

    add_heading1(doc, "2. Baldrige 卓越绩效模型")
    add_paragraph(doc,
        "Baldrige（马尔科姆·波多里奇国家质量奖）是美国国家质量奖，"
        "其评审框架被广泛用于衡量组织的整体卓越绩效。")

    add_heading2(doc, "2.1 Baldrige 七大类")
    add_table(doc,
        ["类别", "名称", "权重", "关键内容"],
        [
            ["1", "领导力", "120分", "高层领导、治理和社会责任"],
            ["2", "战略", "85分", "战略规划和部署"],
            ["3", "顾客", "85分", "顾客洞察、参与和满意度"],
            ["4", "测量、分析和知识管理", "90分", "数据驱动决策和知识管理"],
            ["5", "劳动力", "85分", "员工参与、发展和工作环境"],
            ["6", "运营", "85分", "运营流程设计和管理"],
            ["7", "结果", "450分", "所有领域的绩效结果"],
        ])
    add_paragraph(doc, "总分 1000 分。典型获奖企业得分在 700-800 分以上。")

    add_page_break(doc)

    add_heading1(doc, "3. EFQM 模型")
    add_paragraph(doc,
        "EFQM（European Foundation for Quality Management，欧洲质量管理基金会）"
        "模型是欧洲最广泛使用的卓越运营框架。2020年发布了新版 EFQM Model。")

    add_heading2(doc, "3.1 EFQM 九大要素（新版）")
    add_table(doc,
        ["要素", "名称", "内容"],
        [
            ["1", "Direction（方向）", "组织的使命、愿景、战略方向"],
            ["2", "Execution（执行）", "战略的实施和流程管理"],
            ["3", "Results（成果）", "利益相关方的成果"],
            ["4", "Learn & Improve（学习和改进）", "从结果中学习并持续改进"],
            ["5", "Leadership（领导力）", "高层的愿景和引领"],
            ["6", "People（人才）", "员工的发展和参与"],
            ["7", "Strategy & Execution（战略与执行）", "战略制定与落地"],
            ["8", "Partnerships & Resources（伙伴与资源）", "外部合作和资源管理"],
            ["9", "Innovation & Learning（创新与学习）", "组织的创新能力"],
        ])
    add_paragraph(doc,
        "EFQM 模型采用 RADAR 逻辑进行评估：Results（成果）- Approach（方法）- "
        "Deploy（部署）- Assess & Refine（评估和改进）。")

    add_page_break(doc)

    add_heading1(doc, "4. 精益成熟度整合模型")
    add_paragraph(doc,
        "综合以上各模型，可以构建一个精益成熟度整合模型，用于评估和规划精益转型路径。")

    add_table(doc,
        ["成熟度等级", "特征", "工具水平", "文化水平"],
        [
            ["Level 1 初始级", "精益活动零散，无系统性", "了解基本概念", "被动响应"],
            ["Level 2 规范级", "核心工具推行（5S、标准作业）", "应用基础工具", "开始参与"],
            ["Level 3 整合级", "多工具协同，价值流优化", "熟练运用多种工具", "主动改善"],
            ["Level 4 优化级", "跨部门协同，系统性改善", "定制化方法论", "全员参与"],
            ["Level 5 卓越级", "持续自我进化，行业标杆", "创新方法论", "改善是文化"],
        ])

    add_page_break(doc)

    add_heading1(doc, "5. 卓越运营评估方法")
    add_paragraph(doc, "常用评估方法：")
    add_bullet(doc, "内部审核：由内部评估团队按照模型标准进行系统审核")
    add_bullet(doc, "外部评估：由第三方机构或咨询公司进行独立评估")
    add_bullet(doc, "自评问卷：组织成员自行评分，识别差距")
    add_bullet(doc, "标杆对比（Benchmarking）：与行业最佳实践对比")
    add_bullet(doc, "Gemba Walk（现场走查）：到现场观察实际运营状况")
    add_bullet(doc, "成熟度评估矩阵：按维度和等级进行系统评估")

    add_page_break(doc)

    add_heading1(doc, "6. 行业标杆案例")
    add_heading2(doc, "6.1 Toyota（丰田）")
    add_paragraph(doc,
        "丰田是 TPS 和精益生产的标杆。其卓越运营体现在：全员改善（每年超过100万条建议）、"
        "极低的库存水平、高 OEE、快速交付。丰田模式已成为全球制造业的典范。")

    add_heading2(doc, "6.2 Danaher（丹纳赫）")
    add_paragraph(doc,
        "丹纳赫是全球运营卓越的典范，通过 DBS（Danaher Business System）"
        "系统性地实施精益生产。丹纳赫在50年间进行了400+次收购，每次都能将 DBS"
        "快速复制到新收购的企业中，实现运营改善和价值创造。")

    add_heading2(doc, "6.3 Fastenal（快扣）")
    add_paragraph(doc,
        "Fastenal 是紧固件分销行业的标杆。其 VMI（供应商管理库存）模式和"
        "自动化分销系统体现了卓越运营的理念。通过精益物流和数字化技术，"
        "Fastenal 实现了高效的紧固件分销服务。")

    add_page_break(doc)

    add_heading1(doc, "7. 从精益到卓越的路径")
    add_paragraph(doc,
        "基于各卓越运营模型的共性，以下是制造业从精益到卓越的典型路径：")

    add_heading2(doc, "Phase 1: 基础建设（0-12个月）")
    add_bullet(doc, "5S 现场管理推行")
    add_bullet(doc, "标准作业建立")
    add_bullet(doc, "目视管理实施")
    add_bullet(doc, "基础质量工具培训（SPC、检查表）")

    add_heading2(doc, "Phase 2: 系统构建（12-24个月）")
    add_bullet(doc, "TPM 自主维护和计划维护推行")
    add_bullet(doc, "看板和拉动系统实施")
    add_bullet(doc, "价值流图绘制和改善")
    add_bullet(doc, "快速换型（SMED）推广")

    add_heading2(doc, "Phase 3: 能力提升（24-36个月）")
    add_bullet(doc, "精益六西格玛项目推进")
    add_bullet(doc, "供应链精益化")
    add_bullet(doc, "全员改善文化建立")
    add_bullet(doc, "问题解决能力培养（A3、8D）")

    add_heading2(doc, "Phase 4: 卓越运营（36个月+）")
    add_bullet(doc, "跨部门系统优化")
    add_bullet(doc, "对标世界级水平")
    add_bullet(doc, "创新驱动改善")
    add_bullet(doc, "组织文化的根本转变")

    doc.add_paragraph()
    add_paragraph(doc, "文档结束  |  WCM-METH-005  |  版本 1.0", bold=True)

    path = DOCS_DIR / "05-卓越运营模型.docx"
    doc.save(str(path))
    return path


# ===================================================================
# DOCUMENT 6: Methodology Comparison Matrix
# ===================================================================
def create_comparison_doc():
    doc = Document()
    add_title(doc, "方法论对比矩阵")
    add_paragraph(doc, "六大方法论系统对比与整合指南")
    add_paragraph(doc, "文档编号: WCM-METH-006  |  版本: 1.0  |  分类: 对比分析")
    add_page_break(doc)

    add_heading1(doc, "1. 六大方法论对比总表")
    add_paragraph(doc,
        "以下对比表从多个维度系统比较 TPS、Lean、WCM、Six Sigma、Lean Six Sigma 和 TOC"
        "六大制造方法论。")

    add_heading2(doc, "1.1 基本信息对比")
    add_table(doc,
        ["维度", "TPS", "Lean", "WCM", "Six Sigma", "LSS", "TOC"],
        [
            ["起源时间", "1950s", "1990", "1980s", "1987", "2001", "1984"],
            ["创始人", "大野耐一", "Womack等", "Schonberger", "Motorola/Mikel Harry", "Motorola/GE", "Eliyahu Goldratt"],
            ["起源国家", "日本", "美国/日本", "美国/意大利", "美国", "美国", "以色列"],
            ["代表著作", "丰田生产方式", "改变世界的机器", "World-Class Mfg", "六西格玛管理法", "精益六西格玛", "目标(The Goal)"],
        ])

    add_heading2(doc, "1.2 核心理念对比")
    add_table(doc,
        ["维度", "TPS", "Lean", "WCM", "Six Sigma", "LSS", "TOC"],
        [
            ["核心理念", "消除浪费+自働化", "消除浪费、创造价值", "全面整合卓越", "减少变异、零缺陷", "效率+质量整合", "聚焦瓶颈、突破约束"],
            ["主要支柱", "JIT+自働化", "价值流+流动+拉动", "9大支柱", "DMAIC", "DMAIC+Lean工具", "5步聚焦法"],
            ["驱动方式", "文化驱动", "价值流驱动", "支柱推进驱动", "项目+数据驱动", "项目+数据驱动", "约束驱动"],
            ["关注焦点", "系统整体", "浪费和流动", "全面维度", "变异和质量", "效率和质量", "瓶颈和产出"],
        ])

    add_heading2(doc, "1.3 关键工具对比")
    add_table(doc,
        ["维度", "TPS", "Lean", "WCM", "Six Sigma", "LSS", "TOC"],
        [
            ["核心工具", "看板、安灯、自働化", "VSM、5S、SMED", "TPM+Lean+质量", "SPC、DOE、MSA", "全部Lean+SS工具", "DBR、缓冲管理"],
            ["改善方法", "Kaizen、A3", "Kaizen、VSM改善", "支柱改善", "DMAIC、DFSS", "DMAIC", "5步聚焦法"],
            ["排程方法", "看板拉动", "看板拉动", "综合排程", "基于数据", "综合", "DBR"],
            ["评估工具", "VDA6.3", "成熟度评估", "DC/JD/WD", "西格玛水平", "西格玛水平+OEE", "T/I/OE指标"],
        ])

    add_heading2(doc, "1.4 适用场景对比")
    add_table(doc,
        ["维度", "TPS", "Lean", "WCM", "Six Sigma", "LSS", "TOC"],
        [
            ["最佳适用", "重复性制造", "多品种小批量", "大型制造企业", "质量改进项目", "综合改进项目", "有明显瓶颈的系统"],
            ["行业偏好", "汽车、电子", "通用制造", "汽车、食品、紧固件", "半导体、医疗", "通用", "制造、项目管理"],
            ["企业规模", "中大型", "不限", "大型", "中大型", "中大型", "不限"],
            ["实施复杂度", "高", "中", "很高", "中高", "中高", "中"],
        ])

    add_heading2(doc, "1.5 优势与局限性对比")
    add_table(doc,
        ["维度", "TPS", "Lean", "WCM", "Six Sigma", "LSS", "TOC"],
        [
            ["主要优势", "系统完整、经过验证", "工具丰富、易入门", "全面整合、覆盖广", "数据驱动、严谨", "效率+质量双重提升", "聚焦重点、见效快"],
            ["主要局限", "文化移植困难", "对变异控制弱", "实施周期长、复杂", "对速度关注不够", "需要统计人才", "可能忽视非瓶颈改善"],
            ["投资要求", "中等", "低中", "高", "中等（培训）", "中等", "低中"],
            ["见效速度", "中（2-3年）", "快（3-6月）", "慢（3-5年）", "中（6-12月）", "中（6-12月）", "快（3-6月）"],
        ])

    add_page_break(doc)

    add_heading1(doc, "2. 方法论选择指南（决策矩阵）")
    add_paragraph(doc, "根据企业现状和需求，选择最适合的方法论：")

    add_table(doc,
        ["企业现状", "推荐方法论", "理由"],
        [
            ["刚开始精益转型", "Lean（精益生产）", "工具简单、见效快、文化适应性好"],
            ["设备问题突出（故障频繁）", "TPM → WCM", "先解决设备可靠性，再扩展到全面管理"],
            ["质量问题严重（高不良率）", "Six Sigma → LSS", "数据驱动的质量改进"],
            ["有明显生产瓶颈", "TOC", "聚焦约束、快速突破"],
            ["追求全面卓越运营", "WCM 或 Shingo 模型", "全面覆盖、系统整合"],
            ["需要综合改进", "Lean Six Sigma", "效率和质量同时提升"],
            ["大型跨国制造企业", "WCM + LSS + TOC", "多方法论整合应用"],
        ])

    add_page_break(doc)

    add_heading1(doc, "3. 方法论整合建议（推荐组合）")
    add_paragraph(doc,
        "实际应用中，很少有企业只使用一种方法论。以下是经过验证的推荐整合组合：")

    add_heading2(doc, "3.1 推荐组合一：Lean + TOC（快速见效型）")
    add_paragraph(doc, "适用场景：中小型企业，资源有限，需要快速见效。", bold=True)
    add_bullet(doc, "用 TOC 五步法找到系统瓶颈")
    add_bullet(doc, "在瓶颈应用 Lean 工具消除浪费")
    add_bullet(doc, "用看板管理非瓶颈工序的物料流动")
    add_bullet(doc, "优势：聚焦重点、见效快、投资少")

    add_heading2(doc, "3.2 推荐组合二：TPM + Lean + Six Sigma（全面改进型）")
    add_paragraph(doc, "适用场景：中大型制造企业，有充足的资源和时间。", bold=True)
    add_bullet(doc, "TPM 保障设备可靠性（OEE > 85%）")
    add_bullet(doc, "Lean 消除流程浪费、加快流动")
    add_bullet(doc, "Six Sigma 用统计方法解决变异问题")
    add_bullet(doc, "优势：全面覆盖、系统性改进")

    add_heading2(doc, "3.3 推荐组合三：WCM 框架下的综合应用")
    add_paragraph(doc, "适用场景：大型制造企业，追求世界级水平。", bold=True)
    add_bullet(doc, "以 WCM 9大支柱为框架")
    add_bullet(doc, "LM 支柱用 Lean 工具")
    add_bullet(doc, "QM 支柱用 Six Sigma 工具")
    add_bullet(doc, "AM/PM 支柱用 TPM 方法")
    add_bullet(doc, "瓶颈管理用 TOC 方法")
    add_bullet(doc, "优势：全面整合、有评估体系、可持续推进")

    add_page_break(doc)

    add_heading1(doc, "4. 紧固件行业推荐路径")
    add_paragraph(doc,
        "基于紧固件行业特点（多品种、大批量、连续制造、设备密集、质量要求高），"
        "推荐以下实施路径：")

    add_heading2(doc, "Phase 1: 基础期（0-6个月）-- Lean 基础工具")
    add_bullet(doc, "5S 现场管理：冷镦车间、搓丝车间、热处理车间、成品仓库")
    add_bullet(doc, "标准作业：关键工序建立标准作业指导书")
    add_bullet(doc, "目视管理：在制品标识、设备状态看板、质量控制图")
    add_bullet(doc, "基础培训：全员精益意识培训")

    add_heading2(doc, "Phase 2: 设备保障期（6-12个月）-- TPM 自主维护")
    add_bullet(doc, "自主维护：冷镦机、搓丝机、热处理炉的日常点检")
    add_bullet(doc, "计划维护：关键设备预防性维护计划")
    add_bullet(doc, "OEE 监控：建立 OEE 数据采集和分析系统")
    add_bullet(doc, "目标：OEE 从当前水平提升至 75%")

    add_heading2(doc, "Phase 3: 流动优化期（12-18个月）-- Lean 流程优化")
    add_bullet(doc, "价值流图：绘制从原材料到成品的全流程 VSM")
    add_bullet(doc, "瓶颈识别：用 TOC 方法找到系统瓶颈（通常是热处理）")
    add_bullet(doc, "SMED：冷镦机快速换型，缩短换型时间 50%")
    add_bullet(doc, "看板拉动：建立成品仓库向生产线的看板信号")

    add_heading2(doc, "Phase 4: 质量提升期（18-24个月）-- Six Sigma 质量改进")
    add_bullet(doc, "SPC：关键尺寸（螺纹中径、头部高度等）实施统计过程控制")
    add_bullet(doc, "Cp/Cpk 目标：核心尺寸 Cp >= 1.33, Cpk >= 1.0")
    add_bullet(doc, "绿带培养：培养 3-5 名绿带，领导质量改进项目")
    add_bullet(doc, "FMEA：冷镦、热处理等关键工序进行失效模式分析")

    add_heading2(doc, "Phase 5: 卓越运营期（24个月+）-- WCM 整合")
    add_bullet(doc, "WCM 评估：按照 DC/JD/WD 标准进行内部评估")
    add_bullet(doc, "跨部门协同：打破部门壁垒，优化端到端价值流")
    add_bullet(doc, "标杆对比：与行业标杆企业对标，找差距")
    add_bullet(doc, "文化建设：将改善融入日常工作的习惯和文化")

    doc.add_paragraph()
    add_paragraph(doc, "文档结束  |  WCM-METH-006  |  版本 1.0", bold=True)

    path = DOCS_DIR / "06-方法论对比矩阵.docx"
    doc.save(str(path))
    return path


# ===================================================================
# PPT 7: WCM Training Slides
# ===================================================================
def create_wcm_ppt():
    prs = Presentation()
    prs.slide_width = PInches(13.33)
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE, 7.5)
    ppt_add_title_text(slide, "WCM 世界级制造培训", top=2.0, font_size=44, bold=True)
    ppt_add_title_text(slide, "World-Class Manufacturing Training", top=3.2,
                       font_size=24, color=PRGBColor(0xA0, 0xC0, 0xE0))
    ppt_add_title_text(slide, "精益知识库 | 方法论培训系列", top=4.8,
                       font_size=18, color=PRGBColor(0xC0, 0xD0, 0xE8))

    # --- Slide 2: WCM Overview ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "WCM 概述与愿景", top=0.2)
    lines = [
        "WCM（World-Class Manufacturing）= 世界级制造",
        "",
        "愿景：通过全面整合制造系统的各个维度，达到世界级的运营水平",
        "",
        "核心理念：",
        "  - 零损失（Zero Loss）：零缺陷、零停机、零事故、零库存浪费",
        "  - 全员参与（Total Employee Involvement）",
        "  - 全系统整合（Total System Integration）",
        "  - 持续改善（Continuous Improvement）",
        "  - 客户导向（Customer Focus）",
        "",
        "WCM = Lean + TPM + 质量 + 安全 + 环境 + 成本 + 交期 + 组织效率",
    ]
    ppt_add_content_box(slide, lines, font_size=16)

    # --- Slide 3: 9 Pillars Overview ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "WCM 九大支柱总览", top=0.2)
    ppt_add_table_slide(slide,
        ["支柱", "全称", "核心目标"],
        [
            ["AM", "自主维护", "设备自主管理，零故障"],
            ["PM", "计划维护", "预防性/预测性维护"],
            ["QM", "质量管理", "零缺陷制造"],
            ["DM", "交期管理", "准时交货，缩短交期"],
            ["EM", "环境管理", "绿色制造，合规"],
            ["LM", "精益管理", "消除浪费，持续改善"],
            ["SM", "安全管理", "零事故工作环境"],
            ["OE", "组织效率", "人才发展与组织能力"],
            ["CM", "成本管理", "全员成本意识，持续降本"],
        ])

    # --- Slides 4-12: Each pillar ---
    pillars = [
        ("AM 自主维护", "Autonomous Maintenance",
         ["目标：实现\"我的设备我维护\"",
          "自主维护 7 步法：",
          "  Step 1: 初始清扫 → 发现异常点",
          "  Step 2: 对策与改善 → 修复异常",
          "  Step 3: 清扫润滑基准 → 制定标准",
          "  Step 4: 总点检 → 设备全面教育",
          "  Step 5: 自主点检 → 操作员独立点检",
          "  Step 6: 工作场所标准化 → 5S整合",
          "  Step 7: 自主管理 → 持续自主管理",
          "关键指标：OEE、MTBF、MTTR"]),
        ("PM 计划维护", "Planned Maintenance",
         ["目标：零故障、零停机",
          "维护策略：",
          "  - 预防性维护（PM）：按周期定期维护",
          "  - 预测性维护（PdM）：基于状态监控",
          "  - 改善性维护（CM）：根本性改善",
          "关键活动：",
          "  - 维护计划编制（年度/月度）",
          "  - 振动分析、油液分析、红外热成像",
          "  - 备件管理优化",
          "关键指标：MTBF、MTTR、维护成本率"]),
        ("QM 质量管理", "Quality Management",
         ["目标：零缺陷制造",
          "核心工具：",
          "  - SPC 统计过程控制：控制图监控",
          "  - 过程能力分析：Cp >= 1.33, Cpk >= 1.0",
          "  - 防错（Poka-Yoke）：防止人为失误",
          "  - 8D 问题解决：系统性分析根因",
          "  - 质量功能展开（QFD）：客户需求转化",
          "质量理念：质量是制造出来的，不是检验出来的",
          "关键指标：Cpk、PPM不良率、客户投诉数"]),
        ("DM 交期管理", "Delivery Management",
         ["目标：准时交货，缩短前置时间",
          "关键活动：",
          "  - 价值流图分析（VSM）",
          "  - 前置时间压缩",
          "  - 产线平衡与瓶颈消除",
          "  - 交期承诺管理",
          "  - 订单跟踪与可视化",
          "核心指标：",
          "  - OTD 准时交货率（目标 > 98%）",
          "  - Lead Time 前置时间",
          "  - 订单完成率"]),
        ("EM 环境管理", "Environment Management",
         ["目标：绿色制造，环境合规",
          "关键活动：",
          "  - 废弃物减量化和资源化",
          "  - 能源管理：降低单位产品能耗",
          "  - 有害物质管控（RoHS、REACH）",
          "  - 碳排放管理：碳足迹核算",
          "  - 环境管理体系（ISO 14001）",
          "与法规对接：",
          "  - 国家环保法规",
          "  - 行业排放标准",
          "  - 客户环保要求"]),
        ("LM 精益管理", "Lean Management",
         ["目标：消除浪费，持续改善",
          "核心工具：",
          "  - 5S + 目视管理",
          "  - 标准作业",
          "  - 快速换型（SMED）",
          "  - 看板与拉动系统",
          "  - 价值流管理",
          "改善机制：",
          "  - 日常改善（Daily Kaizen）",
          "  - 课题改善（A3 Problem Solving）",
          "  - 突破性改善（Kaikaku）",
          "  - 改善提案制度"]),
        ("SM 安全管理", "Safety Management",
         ["目标：零事故工作环境",
          "关键活动：",
          "  - 危险源辨识与风险评估（HIRA）",
          "  - 作业安全标准化",
          "  - 安全行为观察（BBS）",
          "  - 应急响应预案",
          "  - 安全培训与意识提升",
          "紧固件行业重点：",
          "  - 冲压设备安全防护",
          "  - 化学品管理（表面处理）",
          "  - 噪音防护",
          "  - 热处理高温防护"]),
        ("OE 组织效率", "Organization Effectiveness",
         ["目标：人才发展与组织能力提升",
          "关键活动：",
          "  - 多能工培养：扩大技能矩阵",
          "  - 团队管理：班组建设和改善活动",
          "  - 绩效管理：目标分解与追踪",
          "  - 知识管理：隐性知识显性化",
          "  - 员工参与：提案制度和改善小组",
          "衡量标准：",
          "  - 人均产出",
          "  - 员工满意度",
          "  - 技能覆盖率",
          "  - 改善提案参与率"]),
        ("CM 成本管理", "Cost Management",
         ["目标：全员成本意识，持续降本",
          "关键活动：",
          "  - 成本可视化：将成本信息传递到一线",
          "  - 损失地图（Loss Map）：量化各种损失",
          "  - 成本改善课题：专项降本活动",
          "  - 价值工程（VA/VE）：设计端降本",
          "  - 采购成本管理：供应链协同",
          "紧固件行业成本重点：",
          "  - 材料利用率（钢材利用率）",
          "  - 能耗成本（电力、天然气）",
          "  - 换型损失（停机时间）",
          "  - 不良品成本（返工、报废）"]),
    ]
    for name, eng_name, content in pillars:
        slide = prs.slides.add_slide(blank_layout)
        ppt_add_header_bar(slide, DARK_BLUE)
        ppt_add_title_text(slide, f"{name}（{eng_name}）", top=0.2)
        ppt_add_content_box(slide, content, font_size=15, top=1.3, height=5.8)

    # --- Slide 13: Implementation Roadmap ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "WCM 实施路线图", top=0.2)
    ppt_add_table_slide(slide,
        ["阶段", "周期", "核心活动", "里程碑"],
        [
            ["导入期", "3-6月", "高层承诺、组织架构、培训、现状评估", "WCM 推进组织成立"],
            ["基础期", "6-12月", "5S、设备维护、标准作业", "OEE > 65%"],
            ["发展期", "12-24月", "9大支柱全面推行、改善活跃", "OEE > 75%"],
            ["成熟期", "24-36月", "系统整合、跨部门协同", "达到 JD 等级"],
            ["卓越期", "36月+", "持续进化、行业标杆", "达到 WD 等级"],
        ])

    # --- Slide 14: Assessment ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "WCM 评估标准", top=0.2)
    ppt_add_table_slide(slide,
        ["评估维度", "DC 等级", "JD 等级", "WD 等级"],
        [
            ["组织与人员", "基本架构建立", "职责明确、积极参与", "卓越领导力和人才发展"],
            ["自愿改善活动", "开始开展", "制度化、常态化", "全员自主改善"],
            ["信息交换", "基本数据收集", "系统化信息共享", "数据驱动决策"],
            ["过程/系统/技术", "标准建立", "系统运行有效", "创新和持续优化"],
            ["结果", "初步改善", "显著改善", "世界级水平"],
        ])

    # --- Slide 15: Relationship ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "WCM 与 Lean / TPM 的关系", top=0.2)
    ppt_add_table_slide(slide,
        ["维度", "WCM", "Lean", "TPM"],
        [
            ["范围", "全面（9大支柱）", "价值流优化", "设备维护"],
            ["核心", "整合所有维度", "消除浪费", "设备可靠性"],
            ["关系", "WCM包含Lean和TPM", "WCM的核心支柱之一", "WCM的AM+PM支柱"],
            ["适用", "大型制造企业", "通用", "设备密集型"],
        ])

    # --- Slide 16: Fastener Industry ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "紧固件行业应用案例", top=0.2)
    lines = [
        "WCM 在紧固件行业的全面应用：",
        "",
        "AM 自主维护：冷镦机、搓丝机、热处理炉日常点检",
        "PM 计划维护：关键设备预防性维护，OEE提升至80%+",
        "QM 质量管理：螺纹通止规、SPC控制图、Cpk提升",
        "DM 交期管理：从冷镦到成品的全流程前置时间压缩",
        "EM 环境管理：废水处理、粉尘收集、能耗优化",
        "LM 精益管理：5S推行、快速换模、看板系统",
        "SM 安全管理：冲压安全防护、噪音控制、化学品管理",
        "OE 组织效率：多能工培养、改善提案制度",
        "CM 成本管理：材料利用率提升、换型损失减少",
        "",
        "预期效果：OEE从60%提升至80%，交期缩短30%，不良率降低60%",
    ]
    ppt_add_content_box(slide, lines, font_size=14, top=1.3, height=5.8)

    # --- Slide 17: Summary ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_BLUE)
    ppt_add_title_text(slide, "总结与行动计划", top=0.2)
    lines = [
        "WCM 核心要点回顾：",
        "",
        "1. WCM 是全面整合的制造管理体系（9大支柱）",
        "2. 以 Lean + TPM 为核心，覆盖质量、安全、环境、成本等维度",
        "3. 采用 DC → JD → WD 分级评估，持续提升",
        "4. 紧固件行业全面适用，每个支柱都有具体应用场景",
        "",
        "行动计划：",
        "  [ ] 第1周：完成 WCM 现状评估（每个支柱打分）",
        "  [ ] 第2周：制定各支柱改善目标和行动计划",
        "  [ ] 第1月：启动 5S 和自主维护（AM+LM）",
        "  [ ] 第3月：引入 OEE 监控和计划维护（PM）",
        "  [ ] 第6月：第一轮 WCM 评估，识别差距",
        "  [ ] 第12月：目标达到 DC 等级",
    ]
    ppt_add_content_box(slide, lines, font_size=15, top=1.3, height=5.8)

    path = PPT_DIR / "05-WCM培训课件.pptx"
    prs.save(str(path))
    return path


# ===================================================================
# PPT 8: Lean Six Sigma Training Slides
# ===================================================================
def create_lss_ppt():
    prs = Presentation()
    prs.slide_width = PInches(13.33)
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN, 7.5)
    ppt_add_title_text(slide, "精益六西格玛培训", top=2.0, font_size=44, bold=True)
    ppt_add_title_text(slide, "Lean Six Sigma Training", top=3.2,
                       font_size=24, color=PRGBColor(0x6E, 0xE0, 0xC0))
    ppt_add_title_text(slide, "精益知识库 | 方法论培训系列", top=4.8,
                       font_size=18, color=PRGBColor(0xA0, 0xD0, 0xB8))

    # --- Slide 2: Overview ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "Lean + Six Sigma 整合概述", top=0.2)
    lines = [
        "精益六西格玛 = Lean（速度和流动）+ Six Sigma（变异和质量）",
        "",
        "Lean 的核心：",
        "  - 消除七大浪费（过量生产、等待、运输、过度加工、库存、动作、不良）",
        "  - 加快价值流动",
        "  - 工具：VSM、5S、看板、SMED、标准作业",
        "",
        "Six Sigma 的核心：",
        "  - 减少过程变异",
        "  - 追求零缺陷（3.4 DPMO）",
        "  - 工具：SPC、DOE、MSA、假设检验、FMEA",
        "",
        "整合价值：Lean 加速流程 + Six Sigma 减少变异 = 更快、更稳定、更高质量",
    ]
    ppt_add_content_box(slide, lines, font_size=16, top=1.3, height=5.8)

    # --- Slide 3: DMAIC Overview ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "DMAIC 方法论总览", top=0.2)
    ppt_add_table_slide(slide,
        ["阶段", "英文", "目标", "核心工具", "输出"],
        [
            ["定义", "Define", "明确问题和目标", "项目章程、SIPOC、CTQ", "项目章程"],
            ["测量", "Measure", "建立数据基线", "MSA、VSM、过程能力", "基线数据"],
            ["分析", "Analyze", "识别根本原因", "鱼骨图、5Why、假设检验", "根本原因清单"],
            ["改进", "Improve", "制定和实施方案", "DOE、防错、SMED", "改进方案"],
            ["控制", "Control", "维持改进成果", "SPC、控制计划、标准作业", "控制计划"],
        ])

    # --- Slides 4-8: DMAIC 5 phases ---
    dmaic_phases = [
        ("Define 定义", DARK_GREEN, [
            "目标：明确问题、目标和范围",
            "",
            "核心活动：",
            "  - 项目章程（Project Charter）",
            "     问题陈述、目标、范围、时间表、团队",
            "  - SIPOC 分析",
            "     Supplier → Input → Process → Output → Customer",
            "  - 客户之声（VOC）收集",
            "     调查、访谈、投诉数据",
            "  - 关键质量特性（CTQ）",
            "     将 VOC 转化为可测量的指标",
            "",
            "紧固件案例：定义\"降低M10螺栓冷镦不良率\"项目",
        ]),
        ("Measure 测量", DARK_GREEN, [
            "目标：收集数据，建立当前基线",
            "",
            "核心活动：",
            "  - 数据收集计划",
            "     测量对象、方法、频率、样本量",
            "  - 测量系统分析（MSA / Gage R&R）",
            "     验证测量系统的准确性和重复性",
            "  - 过程能力基线",
            "     当前 Cp、Cpk、PPM",
            "  - 价值流图（VSM）",
            "     量化时间和库存",
            "",
            "紧固件案例：当前Cpk=0.89，PPM=28,000",
        ]),
        ("Analyze 分析", DARK_GREEN, [
            "目标：识别根本原因和关键影响因素",
            "",
            "核心活动：",
            "  - 鱼骨图（因果图）",
            "     人、机、料、法、环、测",
            "  - 5 Why 分析",
            "     追溯根本原因",
            "  - 假设检验",
            "     t检验、ANOVA、卡方检验",
            "  - 回归分析",
            "     X 和 Y 的数学关系",
            "  - 帕累托分析（80/20法则）",
            "",
            "紧固件案例：根因为材料硬度波动(45%)、模具温度(25%)",
        ]),
        ("Improve 改进", DARK_GREEN, [
            "目标：制定和实施改进方案",
            "",
            "核心活动：",
            "  - 头脑风暴和创意生成",
            "  - 实验设计（DOE）",
            "     系统优化关键参数",
            "  - 防错设计（Poka-Yoke）",
            "     从设计上防止错误",
            "  - 精益工具应用",
            "     5S、SMED、看板等",
            "  - 试运行（Pilot）",
            "     小范围验证效果",
            "",
            "紧固件案例：DOE优化模具温度和润滑参数",
        ]),
        ("Control 控制", DARK_GREEN, [
            "目标：维持改进成果，防止回退",
            "",
            "核心活动：",
            "  - 控制计划（Control Plan）",
            "     监控方法、频率、响应措施",
            "  - SPC 控制图",
            "     实时监控过程状态",
            "  - 标准作业文件更新",
            "     将改进固化为标准",
            "  - 培训和交接",
            "     确保相关人员掌握新流程",
            "  - 项目成果报告",
            "     总结收益和经验教训",
            "",
            "紧固件案例：不良率从2.8%降至0.3%，年节约120万元",
        ]),
    ]
    for title, color, content in dmaic_phases:
        slide = prs.slides.add_slide(blank_layout)
        ppt_add_header_bar(slide, color)
        ppt_add_title_text(slide, title, top=0.2)
        ppt_add_content_box(slide, content, font_size=15, top=1.3, height=5.8)

    # --- Slide 9: Tool Box ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "关键工具箱", top=0.2)
    ppt_add_table_slide(slide,
        ["工具", "英文", "用途", "DMAIC阶段"],
        [
            ["SIPOC", "Supplier-Input-Process-Output-Customer", "确定流程边界", "Define"],
            ["VSM", "Value Stream Mapping", "识别浪费和增值", "Measure"],
            ["FMEA", "Failure Mode & Effects Analysis", "风险评估和预防", "Analyze"],
            ["DOE", "Design of Experiments", "参数优化", "Improve"],
            ["SPC", "Statistical Process Control", "过程监控", "Control"],
            ["MSA", "Measurement System Analysis", "验证测量系统", "Measure"],
            ["Cp/Cpk", "Process Capability", "过程能力评估", "Measure/Control"],
        ])

    # --- Slide 10: Sigma Levels ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "西格玛水平与 DPMO", top=0.2)
    ppt_add_table_slide(slide,
        ["西格玛水平", "DPMO", "合格率", "含义"],
        [
            ["1 sigma", "691,462", "30.85%", "过程极不稳定"],
            ["2 sigma", "308,538", "69.15%", "过程不稳定"],
            ["3 sigma", "66,807", "93.32%", "一般制造水平"],
            ["4 sigma", "6,210", "99.38%", "较好水平（紧固件目标）"],
            ["5 sigma", "233", "99.977%", "优秀水平"],
            ["6 sigma", "3.4", "99.99966%", "世界级水平"],
        ])

    # --- Slide 11: Roles ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "角色体系（GB/BB/MBB）", top=0.2)
    ppt_add_table_slide(slide,
        ["角色", "英文", "职责", "培训", "项目要求"],
        [
            ["黄带", "Yellow Belt", "支持项目团队", "1-2天", "无"],
            ["绿带", "Green Belt", "领导改进项目", "2-3周", "完成1-2个项目"],
            ["黑带", "Black Belt", "全职改进专家", "4-6周", "完成4-6个项目"],
            ["大黑带", "MBB", "战略规划、培训", "持续", "领导组织级推行"],
        ])

    # --- Slide 12: Fastener Case ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "紧固件行业应用案例", top=0.2)
    lines = [
        "案例：降低六角螺母冷镦工序不良率",
        "",
        "项目背景：",
        "  月产量500万件，不良率2.8%（PPM=28,000）",
        "  主要缺陷：头部裂纹(45%) + 对角尺寸超差(30%)",
        "  目标：不良率降至0.5%",
        "",
        "DMAIC 过程：",
        "  Define：项目章程、SIPOC分析",
        "  Measure：30天数据收集，Cpk=0.89",
        "  Analyze：根因为材料硬度波动(45%)、模具温度(25%)",
        "  Improve：DOE优化参数，建立原材料检验标准",
        "  Control：SPC监控，标准化作业，人员培训",
        "",
        "项目成果：",
        "  不良率：2.8% → 0.3%（降低89%）",
        "  Cpk：0.89 → 1.52",
        "  年节约成本：约120万元",
    ]
    ppt_add_content_box(slide, lines, font_size=14, top=1.3, height=5.8)

    # --- Slide 13: Relationship ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "Lean Six Sigma 与其他方法论关系", top=0.2)
    ppt_add_table_slide(slide,
        ["维度", "LSS", "TPS", "WCM", "TOC"],
        [
            ["方法论", "DMAIC数据驱动", "两大支柱+14原则", "9大支柱全面整合", "5步聚焦法"],
            ["改进方式", "项目制", "日常改善+课题", "支柱推进", "约束驱动"],
            ["工具偏好", "统计工具为主", "看板+自働化", "综合工具", "DBR+缓冲管理"],
            ["人才体系", "GB/BB/MBB", "多能工/TWI", "内部评估员", "无特殊角色"],
            ["整合建议", "可作为WCM/TPS的问题解决方法论", "", "", ""],
        ])

    # --- Slide 14: Summary ---
    slide = prs.slides.add_slide(blank_layout)
    ppt_add_header_bar(slide, DARK_GREEN)
    ppt_add_title_text(slide, "总结与下一步", top=0.2)
    lines = [
        "Lean Six Sigma 核心要点回顾：",
        "",
        "1. LSS = Lean（速度）+ Six Sigma（质量）的完美结合",
        "2. DMAIC 是结构化问题解决的核心框架",
        "3. 数据驱动、统计严谨、工具丰富",
        "4. GB/BB/MBB 角色体系支撑持续改进",
        "5. 紧固件行业全面适用",
        "",
        "下一步行动计划：",
        "  [ ] 第1周：选定1个试点项目（建议从质量改进开始）",
        "  [ ] 第2周：启动绿带培训（选3-5名骨干）",
        "  [ ] 第1月：完成试点项目 Define + Measure 阶段",
        "  [ ] 第3月：完成试点项目 Analyze + Improve 阶段",
        "  [ ] 第6月：完成试点项目 Control + 效果验证",
        "  [ ] 第12月：推广至3-5个改进项目",
    ]
    ppt_add_content_box(slide, lines, font_size=15, top=1.3, height=5.8)

    path = PPT_DIR / "06-精益六西格玛培训课件.pptx"
    prs.save(str(path))
    return path


# ===================================================================
# MAIN
# ===================================================================
def main():
    # Ensure directories exist
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    PPT_DIR.mkdir(parents=True, exist_ok=True)

    results = []
    generators = [
        ("01-TPS丰田生产方式详解.docx", create_tps_doc),
        ("02-WCM世界级制造.docx", create_wcm_doc),
        ("03-精益六西格玛.docx", create_lean_six_sigma_doc),
        ("04-TOC约束理论.docx", create_toc_doc),
        ("05-卓越运营模型.docx", create_excellence_model_doc),
        ("06-方法论对比矩阵.docx", create_comparison_doc),
        ("05-WCM培训课件.pptx", create_wcm_ppt),
        ("06-精益六西格玛培训课件.pptx", create_lss_ppt),
    ]

    for name, gen_fn in generators:
        try:
            path = gen_fn()
            file_size = os.path.getsize(path)
            print(f"  [OK] {name}  ->  {path}  ({file_size:,} bytes)")
            results.append((name, True, str(path), file_size))
        except Exception as e:
            print(f"  [FAIL] {name}  ->  Error: {e}")
            results.append((name, False, str(e), 0))

    # Summary
    print("\n" + "=" * 70)
    print("GENERATION SUMMARY")
    print("=" * 70)
    ok_count = sum(1 for _, ok, _, _ in results if ok)
    fail_count = sum(1 for _, ok, _, _ in results if not ok)
    print(f"  Total: {len(results)} | Success: {ok_count} | Failed: {fail_count}")
    for name, ok, path_or_err, size in results:
        status = "OK" if ok else "FAIL"
        detail = f"{path_or_err} ({size:,} bytes)" if ok else path_or_err
        print(f"  [{status}] {name}: {detail}")
    print("=" * 70)

    return 0 if fail_count == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
